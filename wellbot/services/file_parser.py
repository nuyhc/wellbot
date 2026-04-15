"""파일 파싱 서비스.

다양한 파일 형식(PDF/DOCX/XLSX/PPTX/TXT/MD/HWP/이미지)을
텍스트로 변환한다.

파서 종류:
  - LocalParser: 로컬 라이브러리 기반 (pdfplumber, python-docx 등)
  - UpstageParser: Upstage Document Parse API 호출
  - HybridParser: 로컬 우선, 실패 시 Upstage 폴백

팩토리 함수 `get_parser()` 를 통해 설정값에 따라 자동 선택.

PDF 는 Upstage 제약(100p/50MB) 초과 시 자동 분할 후 merge.
그 외 형식은 용량 초과 시 FileTooLargeError 반환 → 사용자 분할 안내.
"""

from __future__ import annotations

import io
import os
import tempfile
from dataclasses import dataclass, field
from pathlib import Path
from typing import Protocol

from wellbot.constants import (
    AUTO_SPLITTABLE_EXTS,
    FILE_PARSER_FALLBACK,
    FILE_PARSER_MODE,
    IMAGE_EXTS,
    LOCAL_SUPPORTED_EXTS,
    SPLIT_SAFETY_PAGES,
    SPLIT_SAFETY_SIZE_MB,
    UPSTAGE_MAX_PAGES,
    UPSTAGE_MAX_SIZE_MB,
    UPSTAGE_SUPPORTED_EXTS,
)


# ── 예외 ──
class ParserError(Exception):
    """파서 공통 예외."""


class UnsupportedFileTypeError(ParserError):
    """현재 모드에서 지원하지 않는 파일 타입."""


class FileTooLargeError(ParserError):
    """파일이 너무 커 자동 분할도 불가능한 경우."""


class ParsingFailedError(ParserError):
    """파싱 중 오류."""


# ── 결과 모델 ──
@dataclass
class ParsedDocument:
    """파싱된 문서 결과.

    Attributes:
        text: 전체 텍스트 (페이지 구분 없음).
        page_count: 페이지/시트/슬라이드 수 (해당 시).
        mime: 판별된 MIME 타입.
        metadata: 추가 메타 (파서가 채움).
    """

    text: str
    page_count: int = 0
    mime: str = ""
    metadata: dict = field(default_factory=dict)


# ── Protocol ──
class DocumentParser(Protocol):
    """문서 파서 인터페이스."""

    def supports(self, ext: str) -> bool:
        """확장자 지원 여부."""
        ...

    def parse(self, file_path: Path) -> ParsedDocument:
        """파일을 파싱해 ParsedDocument 반환."""
        ...


# ── 유틸 ──
def _file_size_mb(path: Path) -> float:
    return path.stat().st_size / (1024 * 1024)


def _ext(path: Path) -> str:
    return path.suffix.lower()


def _count_pdf_pages(pdf_path: Path) -> int:
    """pypdf 로 PDF 페이지 수 조회."""
    from pypdf import PdfReader

    with open(pdf_path, "rb") as f:
        reader = PdfReader(f)
        return len(reader.pages)


# ── 로컬 파서 ──
class LocalParser:
    """로컬 라이브러리 기반 파서."""

    def supports(self, ext: str) -> bool:
        return ext in LOCAL_SUPPORTED_EXTS

    def parse(self, file_path: Path) -> ParsedDocument:
        ext = _ext(file_path)
        if ext == ".pdf":
            return self._parse_pdf(file_path)
        if ext == ".docx":
            return self._parse_docx(file_path)
        if ext == ".xlsx":
            return self._parse_xlsx(file_path)
        if ext == ".pptx":
            return self._parse_pptx(file_path)
        if ext in (".txt", ".md"):
            return self._parse_text(file_path)
        raise UnsupportedFileTypeError(f"LocalParser 는 {ext} 를 지원하지 않습니다.")

    def _parse_pdf(self, path: Path) -> ParsedDocument:
        import pdfplumber

        texts: list[str] = []
        with pdfplumber.open(path) as pdf:
            for page in pdf.pages:
                txt = page.extract_text() or ""
                if txt.strip():
                    texts.append(txt)
            page_count = len(pdf.pages)
        return ParsedDocument(
            text="\n\n".join(texts),
            page_count=page_count,
            mime="application/pdf",
        )

    def _parse_docx(self, path: Path) -> ParsedDocument:
        from docx import Document

        doc = Document(str(path))
        paragraphs = [p.text for p in doc.paragraphs if p.text.strip()]

        # 표도 포함
        for table in doc.tables:
            for row in table.rows:
                cells = [cell.text.strip() for cell in row.cells]
                if any(cells):
                    paragraphs.append(" | ".join(cells))

        return ParsedDocument(
            text="\n".join(paragraphs),
            page_count=0,  # docx 는 페이지 개념 모호
            mime="application/vnd.openxmlformats-officedocument.wordprocessingml.document",
        )

    def _parse_xlsx(self, path: Path) -> ParsedDocument:
        from openpyxl import load_workbook

        wb = load_workbook(str(path), data_only=True, read_only=True)
        sheets_text: list[str] = []
        for sheet in wb.worksheets:
            sheets_text.append(f"## Sheet: {sheet.title}")
            for row in sheet.iter_rows(values_only=True):
                cells = [str(c) if c is not None else "" for c in row]
                if any(cells):
                    sheets_text.append(" | ".join(cells))
        sheet_count = len(wb.sheetnames)
        wb.close()
        return ParsedDocument(
            text="\n".join(sheets_text),
            page_count=sheet_count,
            mime="application/vnd.openxmlformats-officedocument.spreadsheetml.sheet",
        )

    def _parse_pptx(self, path: Path) -> ParsedDocument:
        from pptx import Presentation

        prs = Presentation(str(path))
        slides_text: list[str] = []
        for i, slide in enumerate(prs.slides, start=1):
            slides_text.append(f"## Slide {i}")
            for shape in slide.shapes:
                if hasattr(shape, "text") and shape.text.strip():
                    slides_text.append(shape.text)
        return ParsedDocument(
            text="\n".join(slides_text),
            page_count=len(prs.slides),
            mime="application/vnd.openxmlformats-officedocument.presentationml.presentation",
        )

    def _parse_text(self, path: Path) -> ParsedDocument:
        # UTF-8 우선, 실패 시 cp949 (한국어 Windows 문서)
        for encoding in ("utf-8", "utf-8-sig", "cp949"):
            try:
                text = path.read_text(encoding=encoding)
                mime = "text/markdown" if _ext(path) == ".md" else "text/plain"
                return ParsedDocument(text=text, mime=mime)
            except UnicodeDecodeError:
                continue
        raise ParsingFailedError(f"텍스트 인코딩 감지 실패: {path.name}")


# ── Upstage 파서 ──
class UpstageParser:
    """Upstage Document Parse API 파서.

    제약:
        - 최대 100페이지 / 50MB
        - PDF 는 분할 처리 가능, 그 외 형식은 초과 시 에러
    """

    def __init__(self) -> None:
        self._api_key = os.environ.get("UPSTAGE_API_KEY", "")
        self._api_url = os.environ.get(
            "UPSTAGE_API_URL",
            "https://api.upstage.ai/v1/document-digitization",
        )

    def supports(self, ext: str) -> bool:
        return ext in UPSTAGE_SUPPORTED_EXTS

    def parse(self, file_path: Path) -> ParsedDocument:
        if not self._api_key:
            raise ParserError("UPSTAGE_API_KEY 환경변수가 설정되지 않았습니다.")

        ext = _ext(file_path)
        size_mb = _file_size_mb(file_path)

        # 용량/페이지 초과 → 분할 or 에러
        exceeds_size = size_mb > SPLIT_SAFETY_SIZE_MB
        page_count: int = 0
        if ext == ".pdf":
            page_count = _count_pdf_pages(file_path)
            exceeds_pages = page_count > SPLIT_SAFETY_PAGES
            if exceeds_size or exceeds_pages:
                return self._parse_pdf_split(file_path)
            # 제약 내 → 직접 호출
            result = self._call_api(file_path)
            result.page_count = page_count
            return result

        # 비-PDF: 분할 불가
        if exceeds_size:
            raise FileTooLargeError(
                f"'{ext}' 형식은 자동 분할을 지원하지 않습니다. "
                f"파일 크기 {size_mb:.1f}MB 가 제한 {UPSTAGE_MAX_SIZE_MB}MB 를 "
                f"초과합니다. 파일을 직접 분할하여 재업로드해주세요."
            )

        return self._call_api(file_path)

    def _parse_pdf_split(self, pdf_path: Path) -> ParsedDocument:
        """PDF 를 분할 파싱 후 merge."""
        parts = split_pdf_for_upstage(pdf_path)
        texts: list[str] = []
        total_pages = 0
        try:
            for part_path in parts:
                part_result = self._call_api(part_path)
                texts.append(part_result.text)
                # 각 파트의 페이지 수 합산
                total_pages += _count_pdf_pages(part_path)
        finally:
            # 임시 파일 정리
            for p in parts:
                try:
                    if p != pdf_path:
                        p.unlink(missing_ok=True)
                except Exception:
                    pass

        return ParsedDocument(
            text="\n\n".join(texts),
            page_count=total_pages,
            mime="application/pdf",
            metadata={"split_parts": len(parts)},
        )

    def _call_api(self, file_path: Path) -> ParsedDocument:
        """Upstage API 를 호출한다."""
        import httpx

        headers = {"Authorization": f"Bearer {self._api_key}"}
        with open(file_path, "rb") as f:
            files = {"document": (file_path.name, f, _guess_mime(file_path))}
            data = {
                "model": "document-parse",
                "ocr": "auto",
                "output_formats": '["text"]',
            }
            try:
                response = httpx.post(
                    self._api_url,
                    headers=headers,
                    files=files,
                    data=data,
                    timeout=300.0,
                )
                response.raise_for_status()
                payload = response.json()
            except httpx.HTTPStatusError as e:
                raise ParsingFailedError(
                    f"Upstage API 오류 ({e.response.status_code}): {e.response.text[:200]}"
                ) from e
            except httpx.HTTPError as e:
                raise ParsingFailedError(f"Upstage API 호출 실패: {e}") from e

        # Upstage 응답에서 텍스트 추출
        content = payload.get("content", {})
        text = content.get("text") or ""
        if not text:
            # elements[].content.text concat
            elements = payload.get("elements", []) or []
            parts: list[str] = []
            for el in elements:
                el_text = (el.get("content") or {}).get("text") or ""
                if el_text:
                    parts.append(el_text)
            text = "\n".join(parts)

        usage = payload.get("usage", {}) or {}
        pages = int(usage.get("pages", 0))

        return ParsedDocument(
            text=text,
            page_count=pages,
            mime=_guess_mime(file_path),
            metadata={"upstage_model": payload.get("model", "")},
        )


# ── Hybrid 파서 ──
class HybridParser:
    """로컬 우선 → 실패 시 Upstage 폴백.

    로컬에서 지원하지 않는 형식(HWP 등)은 곧바로 Upstage 사용.
    """

    def __init__(self) -> None:
        self._local = LocalParser()
        self._upstage = UpstageParser()

    def supports(self, ext: str) -> bool:
        return self._local.supports(ext) or self._upstage.supports(ext)

    def parse(self, file_path: Path) -> ParsedDocument:
        ext = _ext(file_path)

        # 로컬 우선
        if self._local.supports(ext):
            try:
                result = self._local.parse(file_path)
                # 빈 결과이고 폴백 허용이면 Upstage 시도
                if (
                    not result.text.strip()
                    and FILE_PARSER_FALLBACK
                    and self._upstage.supports(ext)
                ):
                    return self._upstage.parse(file_path)
                return result
            except Exception:
                if FILE_PARSER_FALLBACK and self._upstage.supports(ext):
                    return self._upstage.parse(file_path)
                raise

        # 로컬 미지원 → Upstage
        if self._upstage.supports(ext):
            return self._upstage.parse(file_path)

        raise UnsupportedFileTypeError(f"지원하지 않는 파일 형식: {ext}")


# ── 팩토리 ──
def get_parser(mode: str | None = None) -> DocumentParser:
    """설정값에 따라 파서를 반환한다."""
    selected = (mode or FILE_PARSER_MODE or "local").lower()
    if selected == "local":
        return LocalParser()
    if selected == "upstage":
        return UpstageParser()
    if selected == "hybrid":
        return HybridParser()
    raise ValueError(f"알 수 없는 FILE_PARSER_MODE: {selected}")


def is_image(file_path: Path | str) -> bool:
    """이미지 파일 여부 (Bedrock Converse image block 으로 직접 전달할 대상)."""
    ext = Path(file_path).suffix.lower()
    return ext in IMAGE_EXTS


# ── PDF 분할 ──
def split_pdf_for_upstage(pdf_path: Path) -> list[Path]:
    """Upstage 제약(페이지/용량)에 맞게 PDF 를 분할한다.

    분할 단계:
        1. SPLIT_SAFETY_PAGES 단위로 페이지 기준 1차 분할
        2. 각 파트가 UPSTAGE_MAX_SIZE_MB 초과 시 재귀적 반분할
        3. 단일 페이지가 50MB 초과 시 FileTooLargeError

    Returns:
        분할된 임시 PDF 파일 경로 목록. 페이지 순으로 정렬됨.
        입력 파일이 제약 내면 [pdf_path] 를 그대로 반환.
    """
    from pypdf import PdfReader, PdfWriter

    total_pages = _count_pdf_pages(pdf_path)
    file_size_mb = _file_size_mb(pdf_path)

    if total_pages <= SPLIT_SAFETY_PAGES and file_size_mb <= SPLIT_SAFETY_SIZE_MB:
        return [pdf_path]

    parts: list[Path] = []
    reader = PdfReader(str(pdf_path))

    # 1차 분할: 페이지 기준
    for start in range(0, total_pages, SPLIT_SAFETY_PAGES):
        end = min(start + SPLIT_SAFETY_PAGES, total_pages)
        part_path = _write_pdf_part(reader, start, end, pdf_path.stem, len(parts))
        parts.append(part_path)

    # 2차 검증: 용량 초과 파트는 반분할
    final_parts: list[Path] = []
    for part in parts:
        final_parts.extend(_ensure_size_limit(part))

    return final_parts


def _write_pdf_part(
    reader,
    start: int,
    end: int,
    stem: str,
    idx: int,
) -> Path:
    """PDF 의 [start, end) 페이지를 임시 파일로 저장."""
    from pypdf import PdfWriter

    writer = PdfWriter()
    for i in range(start, end):
        writer.add_page(reader.pages[i])

    tmp_dir = Path(tempfile.gettempdir()) / "wellbot_pdf_split"
    tmp_dir.mkdir(parents=True, exist_ok=True)
    out_path = tmp_dir / f"{stem}_part{idx:03d}_p{start}-{end - 1}.pdf"
    with open(out_path, "wb") as f:
        writer.write(f)
    return out_path


def _ensure_size_limit(pdf_path: Path) -> list[Path]:
    """파일이 용량 제한 초과 시 재귀적으로 반분할한다.

    단일 페이지가 제한 초과면 FileTooLargeError.
    """
    from pypdf import PdfReader

    size_mb = _file_size_mb(pdf_path)
    if size_mb <= UPSTAGE_MAX_SIZE_MB:
        return [pdf_path]

    reader = PdfReader(str(pdf_path))
    page_count = len(reader.pages)
    if page_count <= 1:
        raise FileTooLargeError(
            f"단일 페이지 크기 {size_mb:.1f}MB 가 제한 "
            f"{UPSTAGE_MAX_SIZE_MB}MB 를 초과합니다: {pdf_path.name}"
        )

    # 반분할
    mid = page_count // 2
    left = _write_pdf_part(reader, 0, mid, pdf_path.stem, 0)
    right = _write_pdf_part(reader, mid, page_count, pdf_path.stem, 1)

    # 원본(임시파일) 정리
    try:
        pdf_path.unlink()
    except Exception:
        pass

    # 재귀 검증
    result: list[Path] = []
    result.extend(_ensure_size_limit(left))
    result.extend(_ensure_size_limit(right))
    return result


# ── MIME 추출 ──
_MIME_MAP = {
    ".pdf": "application/pdf",
    ".docx": "application/vnd.openxmlformats-officedocument.wordprocessingml.document",
    ".xlsx": "application/vnd.openxmlformats-officedocument.spreadsheetml.sheet",
    ".pptx": "application/vnd.openxmlformats-officedocument.presentationml.presentation",
    ".txt": "text/plain",
    ".md": "text/markdown",
    ".hwp": "application/x-hwp",
    ".hwpx": "application/x-hwpx",
    ".png": "image/png",
    ".jpg": "image/jpeg",
    ".jpeg": "image/jpeg",
    ".webp": "image/webp",
    ".tiff": "image/tiff",
    ".bmp": "image/bmp",
}


def _guess_mime(file_path: Path) -> str:
    """확장자로부터 MIME 타입을 추정한다."""
    ext = _ext(file_path)
    return _MIME_MAP.get(ext, "application/octet-stream")


def guess_mime(file_path: Path | str) -> str:
    """외부 공개용 MIME 추정."""
    return _guess_mime(Path(file_path))


def count_pdf_pages(pdf_path: Path | str) -> int:
    """외부 공개용 PDF 페이지 수 조회."""
    return _count_pdf_pages(Path(pdf_path))
