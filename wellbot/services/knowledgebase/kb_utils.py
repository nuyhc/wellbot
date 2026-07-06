"""
kb_utils.py

KB 매니저 공통 인프라.
personal_kb_manager, team_kb_manager 에서 공유하는 상수/클라이언트/함수.

- 설정 상수 (S3 버킷, Lambda ARN 등)
- AWS 클라이언트 (s3, bedrock-agent, s3vectors)
- 파일 크기 검증
- xlsx/csv 분할 업로드
- pptx → json 변환 (Bedrock KB 미지원 형식 전처리)
- KB 생성/조회/Ingestion (kind="personal"|"team" 으로 분기)
"""

import io
import json
import logging
import os
import re
import tempfile
import threading
import time
from concurrent.futures import ThreadPoolExecutor, as_completed
from functools import lru_cache
from pathlib import Path
from typing import Optional

import boto3
import pandas as pd

from wellbot.constants import FILE_PARSER_MODE, KB_MAX_DOCS, PDF_VIA_UPSTAGE
from wellbot.services.files import storage_service
from wellbot.services.knowledgebase.config import get_kb_config

log = logging.getLogger(__name__)


# ──────────────────────────────────────────────
# 식별자
# ──────────────────────────────────────────────
AGNT_ID_KB    = "AgntKnowBase"
KB_INFO_SEP   = "||"

# kind → S3 경로 / KB 이름에 사용되는 base 토큰
_KB_KIND_PREFIX_BASE = {"personal": "users", "team": "teams"}
_KB_KIND_LABEL       = {"personal": "개인",  "team": "팀"}


@lru_cache(maxsize=1)
def env_suffix() -> str:
    """환경 구분 접미사. APP_ENV 값으로 결정.

    dev/prd 가 같은 AWS 계정·버킷을 공유하므로 KB/DS/벡터인덱스 이름과 S3
    prefix 에 이 값을 붙여 충돌을 방지.
    - 미설정 → 'dev'(기본) → '-dev'   : 빠뜨려도 prod 네임스페이스를 오염하지 않음
    - APP_ENV='' / 'prod' / 'prd'      : ''(빈 문자열) → 기존 이름·경로 유지(마이그레이션 불필요)
    - APP_ENV='dev'/'staging'/...      : '-<env>'

    lazy(lru_cache) — init_env() 이후 첫 호출 시 1회 평가.
    """
    env = os.getenv("APP_ENV", "dev").strip().lower()
    return "" if env in ("", "prod", "prd") else f"-{env}"


def kb_base(kind: str) -> str:
    """S3 경로 첫 세그먼트: 'users'/'users-dev'/'teams'/'teams-dev'."""
    return f"{_KB_KIND_PREFIX_BASE[kind]}{env_suffix()}"


def shared_base() -> str:
    """공용 KB S3 경로 첫 세그먼트: 'shared'/'shared-dev'. (개인/팀의 kb_base 와 동일 규칙.)

    공용 KB는 개인/팀과 달리 'shared/{대분류}/raw/{소분류}' 2단계 구조라 kb_base/
    raw_prefix 를 재사용하지 않고 이 base 만 공유한다(목록 조회·업로드·DS 생성 단일 출처).
    """
    return f"shared{env_suffix()}"


# ──────────────────────────────────────────────
# 설정 / 리전 (lazy — import 사이드이펙트 방지)
# ──────────────────────────────────────────────
@lru_cache(maxsize=1)
def _cfg() -> dict:
    """personal_kb 섹션 (팀 KB 도 공유). 최초 호출 시 1회 로드 후 캐싱."""
    return get_kb_config()["personal_kb"]


def _region() -> str:
    return os.getenv("AWS_REGION", "ap-northeast-2")


def get_s3_bucket() -> str:
    """KB 파일 저장 버킷. personal/team 매니저가 업로드/삭제 시 사용."""
    return _cfg()["s3_bucket"]


# ──────────────────────────────────────────────
# AWS 클라이언트 (lazy — 최초 호출 시 1회 생성 후 캐싱)
# ──────────────────────────────────────────────
def _get_s3():
    """S3 클라이언트 — storage_service 의 region 설정 클라이언트 재사용"""
    return storage_service.get_client()


@lru_cache(maxsize=1)
def _get_bedrock_agent():
    return boto3.client("bedrock-agent", region_name=_region())


@lru_cache(maxsize=1)
def _get_s3vectors():
    return boto3.client("s3vectors", region_name=_region())

# ──────────────────────────────────────────────
# 상수
# ──────────────────────────────────────────────
ROWS_PER_SPLIT = 50_000          # xlsx/csv 행 기준 분할 단위

TABULAR_EXTS = {".xlsx", ".csv"}

# Bedrock KB 가 지원하지 않지만, 업로드 시 변환 처리하는 형식
CONVERTIBLE_EXTS = {".pptx"}


def get_originals_prefix(raw_prefix: str) -> str:
    """raw/ prefix 를 originals/ prefix 로 변환.

    pptx 등 변환이 필요한 파일의 원본은 Bedrock 의 inclusionPrefix(raw/) 밖에
    저장해서 ingestion 대상에서 제외. 다운로드와 문서 목록 조회용도로만 사용.

    예: 'users/123/raw/'  → 'users/123/originals/'
        'teams/A1/raw/'   → 'teams/A1/originals/'
    """
    return raw_prefix.rsplit("raw/", 1)[0] + "originals/"


def get_staging_prefix(raw_prefix: str) -> str:
    """raw/ prefix 를 staging/ prefix 로 변환.

    업로드 HTTP 요청은 원본을 staging/ 에만 빠르게 적재하고 즉시 반환하고,
    변환(Upstage 등)·색인은 백그라운드(on_upload_complete)에서 staging/ 의 원본을
    읽어 수행한다 → 다중 PDF 등에서 HTTP 프록시 타임아웃(504)과 분리.
    raw/ 의 형제 위치라 Bedrock inclusionPrefix(raw/) 밖 → 색인 대상에서 제외.

    예: 'users/123/raw/'  → 'users/123/staging/'
        'teams/A1/raw/'   → 'teams/A1/staging/'
    """
    return raw_prefix.rsplit("raw/", 1)[0] + "staging/"

SUPPORTED_EXTENSIONS = {
    ".pdf", ".docx", ".pptx", ".xlsx", ".csv",
    ".md", ".txt", ".json", ".html", ".htm",
}

MAX_FILE_SIZES = {
    ".txt":  30 * 1024 * 1024,   # 30MB
    ".md":   30 * 1024 * 1024,
    ".json": 30 * 1024 * 1024,
    ".csv":  None,               # 분할 업로드 — 제한 없음
    ".xlsx": None,
}
MAX_FILE_SIZE_DEFAULT = 100 * 1024 * 1024  # 100MB


# ──────────────────────────────────────────────
# 파일 크기 검증
# ──────────────────────────────────────────────
def validate_file_size(file_bytes: bytes, filename: str) -> None:
    """
    파일 형식별 크기 제한 검증.
    xlsx/csv 는 분할 업로드로 처리되므로 제한 없음.
    """
    ext = Path(filename).suffix.lower()
    limit = MAX_FILE_SIZES.get(ext, MAX_FILE_SIZE_DEFAULT)
    if limit is None:
        return
    if len(file_bytes) > limit:
        limit_mb = limit // (1024 * 1024)
        actual_mb = len(file_bytes) / (1024 * 1024)
        raise ValueError(
            f"파일 크기 초과: {filename} ({actual_mb:.1f}MB). "
            f"{ext} 파일은 {limit_mb}MB 이하만 업로드 가능합니다."
        )


# ──────────────────────────────────────────────
# 분할본(_partN) ↔ 논리 문서 매핑 (단일 출처)
#   xlsx/csv 는 ROWS_PER_SPLIT 단위로 '{stem}_part{N}.ext'(csv·단일시트) 또는
#   '{stem}_{sheet}_part{N}.ext'(멀티시트) 로 분할 저장된다. 카운트/목록/삭제가
#   이 분할본을 하나의 논리 문서로 일관되게 다루도록 매핑 로직을 여기로 통일.
# ──────────────────────────────────────────────
# 분할본 파일명 매처. base = '_part{N}.ext' 를 떼어낸 부분(멀티시트면 시트 슬러그 포함).
_TABULAR_PART_RE = re.compile(r"^(?P<base>.+?)_part\d+(?P<ext>\.(?:xlsx|csv))$", re.IGNORECASE)


def is_tabular_part(filename: str) -> bool:
    """파일명이 xlsx/csv 분할본('..._part{N}.xlsx|csv')인지 여부."""
    return _TABULAR_PART_RE.match(filename) is not None


def _part_owner(base: str, original_stems: set[str]) -> Optional[str]:
    """분할본 base 의 소유 원본 stem 을 판정.

    base 예: 'data'(csv·단일시트) 또는 'budget_매출'(멀티시트, 시트 슬러그 포함).
    소유 stem = base 와 정확히 같거나 base 의 prefix('{stem}_...')인 원본 stem 중
    **가장 긴 것**. originals/ 의 실제 원본 집합으로 판정하므로 prefix 충돌
    (예: 'report' vs 'report_summary')에서도 올바른 소유자를 고른다.
    """
    candidates = [s for s in original_stems if base == s or base.startswith(s + "_")]
    return max(candidates, key=len) if candidates else None


def _list_original_stems(bucket: str, originals_prefix: str) -> set[str]:
    """originals/ 에 보관된 원본 파일들의 stem 집합 (분할본 소유 판정용)."""
    stems: set[str] = set()
    paginator = _get_s3().get_paginator("list_objects_v2")
    for page in paginator.paginate(Bucket=bucket, Prefix=originals_prefix):
        for obj in page.get("Contents", []) or []:
            stems.add(Path(obj["Key"].split("/")[-1]).stem)
    return stems


def _owned_part_keys(
    bucket: str, prefix: str, target_stem: str, ext: str, original_stems: set[str],
) -> list[str]:
    """prefix(raw/) 아래에서 target_stem(논리 파일) 소유의 분할본 key 목록.

    소유 판정에 target_stem 자신도 후보에 포함 — 최초 업로드(원본 아직 미보관)나
    재업로드 어느 경우든 자기 분할본을 정확히 집어낸다.
    """
    pool = set(original_stems) | {target_stem}
    keys: list[str] = []
    paginator = _get_s3().get_paginator("list_objects_v2")
    for page in paginator.paginate(Bucket=bucket, Prefix=prefix):
        for obj in page.get("Contents", []) or []:
            key = obj["Key"]
            filename = key.split("/")[-1]
            m = _TABULAR_PART_RE.match(filename)
            if not m or m.group("ext").lower() != ext.lower():
                continue
            if _part_owner(m.group("base"), pool) == target_stem:
                keys.append(key)
    return keys


# ──────────────────────────────────────────────
# 분할 업로드 (xlsx / csv)
# ──────────────────────────────────────────────
def cleanup_existing_parts(
    bucket: str, prefix: str, stem: str, ext: str,
) -> None:
    """
    동일 논리 파일의 기존 분할본을 S3 에서 삭제.
    재업로드 시 파트 수/시트 구성이 달라져 오래된 파트가 남는 것을 방지.

    소유 판정을 originals/ 원본 집합 기준으로 하여, 이름이 prefix 로 겹치는
    다른 파일(예: 'report' 재업로드가 'report_summary' 의 분할본을 지우는 것)을 보존.
    """
    original_stems = _list_original_stems(bucket, get_originals_prefix(prefix))
    for key in _owned_part_keys(bucket, prefix, stem, ext, original_stems):
        _get_s3().delete_object(Bucket=bucket, Key=key)


def _safe_sheet_slug(sheet_name: object) -> str:
    """시트명을 S3 키/파일명에 안전한 슬러그로 변환 (한글 보존, 특수문자 → _)."""
    slug = re.sub(r"[^\w가-힣]+", "_", str(sheet_name)).strip("_")
    return slug or "sheet"


def _upload_df_parts(
    bucket: str,
    prefix: str,
    df: "pd.DataFrame",
    part_stem: str,
    ext: str,
    sheet_name: object = None,
) -> list[str]:
    """단일 DataFrame 을 ROWS_PER_SPLIT 행 단위로 분할해 S3 업로드. 반환: URI 목록.

    행이 0개인 (헤더만 있는) 시트는 색인할 내용이 없으므로 스킵.
    xlsx 는 원본 시트명을 파트 내부 시트명으로 보존 → Lambda 가 sheet 메타 태깅.
    """
    uris: list[str] = []
    total_rows = len(df)
    if total_rows == 0:
        return uris
    for i, start in enumerate(range(0, total_rows, ROWS_PER_SPLIT)):
        chunk_df = df.iloc[start:start + ROWS_PER_SPLIT]
        split_filename = f"{part_stem}_part{i + 1}{ext}"
        buf = io.BytesIO()

        if ext == ".csv":
            chunk_df.to_csv(buf, index=False)
        else:
            xl_sheet = (str(sheet_name)[:31] if sheet_name else "") or "Sheet1"
            with pd.ExcelWriter(buf, engine="openpyxl") as writer:
                chunk_df.to_excel(writer, index=False, sheet_name=xl_sheet)

        buf.seek(0)
        key = f"{prefix}{split_filename}"
        _get_s3().put_object(Bucket=bucket, Key=key, Body=buf.read())
        uris.append(f"s3://{bucket}/{key}")
    return uris


def split_and_upload_tabular(
    bucket: str,
    prefix: str,
    file_bytes: bytes,
    filename: str,
) -> list[str]:
    """
    xlsx/csv 를 ROWS_PER_SPLIT 행 단위로 분할해서 S3 에 저장.
    업로드 전 기존 분할 파트를 정리해서 오래된 데이터 잔류를 방지.

    xlsx 는 **모든 시트**를 각각 분할 업로드 (시트명 보존). 멀티시트면 파일명에
    시트 슬러그를 포함('{stem}_{sheet}_part{N}.xlsx'), 단일시트면 기존 명명 유지.
    반환: 업로드된 S3 URI 목록
    """
    ext = Path(filename).suffix.lower()
    stem = Path(filename).stem

    cleanup_existing_parts(bucket, prefix, stem, ext)

    if ext == ".csv":
        try:
            df = pd.read_csv(io.BytesIO(file_bytes))
        except UnicodeDecodeError:
            df = pd.read_csv(io.BytesIO(file_bytes), encoding="cp949")
        return _upload_df_parts(bucket, prefix, df, stem, ext)

    # xlsx: 전체 시트를 dict 로 읽어 시트별 분할 (sheet_name=None → {시트명: df})
    sheets = pd.read_excel(io.BytesIO(file_bytes), sheet_name=None)
    multi = len(sheets) > 1
    uris: list[str] = []
    for sheet_name, df in sheets.items():
        part_stem = f"{stem}_{_safe_sheet_slug(sheet_name)}" if multi else stem
        uris.extend(
            _upload_df_parts(bucket, prefix, df, part_stem, ext, sheet_name=sheet_name)
        )
    return uris


# ──────────────────────────────────────────────
# pptx → json 변환
# ──────────────────────────────────────────────
def pptx_to_dict(file_bytes: bytes) -> dict:
    """pptx 바이트를 슬라이드별 구조화 dict 로 추출 (제목/본문/표/노트).

    convert_pptx_to_json(개인·팀 업로드) 과 scripts/shared_kb_manager(공용 KB)
    가 공유하는 코어 추출 로직. 키는 'slide_{번호}_{제목}', 값은 슬라이드 텍스트.
    """
    from pptx import Presentation

    prs = Presentation(io.BytesIO(file_bytes))
    result: dict = {}

    for slide_num, slide in enumerate(prs.slides, 1):
        slide_title = "No Title"
        if slide.shapes.title and slide.shapes.title.has_text_frame:
            slide_title = slide.shapes.title.text.strip() or "No Title"

        parts = []

        for shape in slide.shapes:
            if shape.has_table:
                table = shape.table
                rows = [[cell.text.strip() for cell in row.cells] for row in table.rows]
                if rows:
                    headers = rows[0]
                    for row in rows[1:]:
                        row_text = " | ".join(
                            f"{headers[i]}: {cell}" for i, cell in enumerate(row)
                            if i < len(headers)
                        )
                        if row_text.strip():
                            parts.append(row_text)
            elif shape.has_text_frame:
                for para in shape.text_frame.paragraphs:
                    text = para.text.strip()
                    if text:
                        parts.append(text)

        if slide.has_notes_slide:
            notes = slide.notes_slide.notes_text_frame.text.strip()
            if notes:
                parts.append(f"[노트] {notes}")

        if parts:
            result[f"slide_{slide_num}_{slide_title}"] = "\n".join(parts)

    return result


def convert_pptx_to_json(file_bytes: bytes, filename: str) -> tuple[bytes, str]:
    """pptx 바이트를 슬라이드별 JSON 으로 변환 (Bedrock KB 미지원 형식 전처리).

    Lambda 의 parse_json 이 처리 가능하도록 업로드 전 변환.

    반환: (json_bytes, 변환된_파일명)
        예: ("report.pptx" → b'{...}', "report_pptx.json")
    """
    result = pptx_to_dict(file_bytes)
    stem = Path(filename).stem
    json_filename = f"{stem}_pptx.json"
    json_bytes = json.dumps(result, ensure_ascii=False, indent=2).encode("utf-8")
    return json_bytes, json_filename


# ──────────────────────────────────────────────
# Upstage Document Parse → markdown 변환 (xlsx / pdf)
# ──────────────────────────────────────────────
def xlsx_via_upstage_enabled() -> bool:
    """xlsx 를 Upstage 로 변환할지 여부 (FILE_PARSER_MODE 가 upstage/hybrid 일 때).

    첨부파일 파서와 동일한 노브(FILE_PARSER_MODE)를 공유한다.
    local 모드면 기존 pandas 행 분할(split_and_upload_tabular)을 그대로 사용.
    """
    return (FILE_PARSER_MODE or "local").lower() in ("upstage", "hybrid")


def pdf_via_upstage_enabled() -> bool:
    """PDF 를 Upstage 로 변환할지 여부 (xlsx 와 독립된 PDF 전용 노브 PDF_VIA_UPSTAGE).

    개인/팀/공용 KB 업로드 전부에 적용. 끄면 원본 PDF 가 그대로 색인되어
    Lambda 의 pdfplumber 커스텀 파싱(parse_pdf)으로 폴백한다.
    """
    return bool(PDF_VIA_UPSTAGE)


def _convert_via_upstage(
    file_bytes: bytes,
    filename: str,
    out_suffix: str,
    with_page_markers: bool = False,
) -> tuple[bytes, str]:
    """파일을 Upstage Document Parse 로 markdown 변환 (xlsx/pdf 공용).

    변환된 markdown 은 Lambda 의 parse_md 가 청킹한다.
    반환: (markdown_bytes, f"{stem}{out_suffix}")
    Upstage 호출 실패/빈 결과 시 예외 전파 (호출자가 폴백 처리).

    with_page_markers=True (PDF 전용) 이면 parsed.pages 로 페이지별 블록 앞에
    `<!--page=N-->` 마커를 삽입한 md 를 만든다 → Lambda parse_md 가 청크 page 메타로 태깅.
    pages 가 비면 parsed.text 로 폴백(마커 없음, 현행 동작).
    """
    # 모듈 레벨 import 사이드이펙트 방지를 위해 지연 import
    from wellbot.services.files.file_parser import UpstageParser

    # 임시 파일을 원본 파일명으로 생성 (Upstage 로그/제목에 임의 임시명이 노출되지 않도록).
    # filename 은 업로드 basename 이라 경로 구분자가 없지만 방어적으로 .name 사용.
    safe_name = Path(filename).name or "upload"
    with tempfile.TemporaryDirectory() as tmpdir:
        tmp_path = Path(tmpdir) / safe_name
        tmp_path.write_bytes(file_bytes)
        parsed = UpstageParser().parse(tmp_path)

    if with_page_markers and parsed.pages:
        md = "\n".join(f"<!--page={pg}-->\n{block}" for pg, block in parsed.pages)
    else:
        md = parsed.text or ""
    if not md.strip():
        raise ValueError("Upstage 변환 결과가 비어 있습니다")
    stem = Path(filename).stem
    return md.encode("utf-8"), f"{stem}{out_suffix}"


def convert_xlsx_to_markdown(file_bytes: bytes, filename: str) -> tuple[bytes, str]:
    """xlsx 를 Upstage 로 markdown 변환 (병합·공백 많은 표를 견고하게 처리).
    예: report.xlsx → report_xlsx.md. 실패 시 예외 전파(호출자가 pandas 분할 폴백)."""
    return _convert_via_upstage(file_bytes, filename, "_xlsx.md")


def convert_pdf_to_markdown(file_bytes: bytes, filename: str) -> tuple[bytes, str]:
    """PDF 를 Upstage 로 markdown 변환 (이미지/스캔 내용까지 OCR+레이아웃 해석).
    예: report.pdf → report_pdf.md. 실패 시 예외 전파(호출자가 원본 PDF 색인 폴백).
    페이지 마커(`<!--page=N-->`)를 삽입해 출처 page 표시가 가능하도록 한다."""
    return _convert_via_upstage(file_bytes, filename, "_pdf.md", with_page_markers=True)


# ──────────────────────────────────────────────
# KB 정보 인코딩 / 디코딩 (DB path_addr 컬럼용)
# ──────────────────────────────────────────────
def encode_kb_info(kb_id: str, data_source_id: str) -> str:
    return f"{kb_id}{KB_INFO_SEP}{data_source_id}"


def decode_kb_info(kb_info: str) -> tuple[str, str]:
    parts = kb_info.split(KB_INFO_SEP, 1)
    if len(parts) != 2:
        raise ValueError(f"잘못된 KB_INFO 형식: {kb_info}")
    return parts[0], parts[1]


# ──────────────────────────────────────────────
# S3 경로 헬퍼
# ──────────────────────────────────────────────
def raw_prefix(kind: str, owner: str) -> str:
    """kind='personal'→'users{env}/{owner}/raw/', kind='team'→'teams{env}/{owner}/raw/'."""
    return f"{kb_base(kind)}/{owner}/raw/"


def processed_prefix(kind: str, owner: str) -> str:
    """processed/ prefix. raw_prefix 와 동일 규칙."""
    return f"{kb_base(kind)}/{owner}/processed/"


# ──────────────────────────────────────────────
# Bedrock KB 생성 / 조회
# ──────────────────────────────────────────────
def create_vector_index(kind: str, owner: str) -> str:
    """S3 Vectors 인덱스 생성. 반환: index ARN."""
    resp = _get_s3vectors().create_index(
        vectorBucketName=_cfg()["s3_vector_bucket"],
        indexName=f"aiinno-bedrock-kb-{kind}-vector-index-{owner.lower()}{env_suffix()}",
        dataType="float32",
        dimension=1024,
        distanceMetric="cosine",
        metadataConfiguration={
            "nonFilterableMetadataKeys": ["AMAZON_BEDROCK_TEXT"]
        },
    )
    return resp["indexArn"]


def create_bedrock_kb(kind: str, owner: str, vector_index_arn: str) -> str:
    """Bedrock Knowledge Base 생성. 반환: knowledgeBaseId."""
    label = _KB_KIND_LABEL[kind]
    resp = _get_bedrock_agent().create_knowledge_base(
        name=f"aiinno-bedrock-kb-{kind}-{owner}{env_suffix()}",
        description=f"{owner}의 {label} Knowledge Base",
        roleArn=_cfg()["kb_role_arn"],
        knowledgeBaseConfiguration={
            "type": "VECTOR",
            "vectorKnowledgeBaseConfiguration": {
                "embeddingModelArn": (
                    f"arn:aws:bedrock:{_region()}::foundation-model/{_cfg()['embedding_model']}"
                ),
            },
        },
        storageConfiguration={
            "type": "S3_VECTORS",
            "s3VectorsConfiguration": {"indexArn": vector_index_arn},
        },
    )
    return resp["knowledgeBase"]["knowledgeBaseId"]

def wait_until_kb_ready(kb_id: str) -> None:
    """KB 가 ACTIVE 가 될 때까지 폴링."""
    cfg = _cfg()
    poll_timeout = cfg.get("kb_poll_timeout", 120)
    poll_interval = cfg.get("kb_poll_interval", 5)
    start = time.time()
    while time.time() - start < poll_timeout:
        resp = _get_bedrock_agent().get_knowledge_base(knowledgeBaseId=kb_id)
        status = resp["knowledgeBase"]["status"]
        if status == "ACTIVE":
            return
        if status == "FAILED":
            raise RuntimeError(f"KB 생성 실패: kb_id={kb_id}")
        time.sleep(poll_interval)
    raise TimeoutError(f"KB ACTIVE 대기 타임아웃: kb_id={kb_id}")


def create_data_source(kind: str, owner: str, kb_id: str) -> str:
    """KB 의 Data Source 생성. 반환: dataSourceId."""
    resp = _get_bedrock_agent().create_data_source(
        knowledgeBaseId=kb_id,
        name=f"aiinno-bedrock-kb-ds-{kind}-{owner}{env_suffix()}",
        dataSourceConfiguration={
            "type": "S3",
            "s3Configuration": {
                "bucketArn": f"arn:aws:s3:::{_cfg()['s3_bucket']}",
                "inclusionPrefixes": [raw_prefix(kind, owner)],
            },
        },
        vectorIngestionConfiguration={
            "chunkingConfiguration": {"chunkingStrategy": "NONE"},
            "customTransformationConfiguration": {
                "intermediateStorage": {
                    "s3Location": {
                        "uri": f"s3://{_cfg()['s3_intermediate_bucket']}/{processed_prefix(kind, owner)}",
                    },
                },
                "transformations": [{
                    "stepToApply": "POST_CHUNKING",
                    "transformationFunction": {
                        "transformationLambdaConfiguration": {"lambdaArn": _cfg()['lambda_arn']},
                    },
                }],
            },
        },
    )
    return resp["dataSource"]["dataSourceId"]


def find_existing_kb(kind: str, owner: str) -> Optional[dict]:
    """Bedrock API 로 이미 생성된 KB 검색 (DB 미등록 상태 방어).

    반환: {"kb_id", "data_source_id"} 또는 None.
    """
    kb_name = f"aiinno-bedrock-kb-{kind}-{owner}{env_suffix()}"
    try:
        paginator = _get_bedrock_agent().get_paginator("list_knowledge_bases")
        for page in paginator.paginate():
            for kb in page.get("knowledgeBaseSummaries", []):
                if kb["name"] == kb_name and kb["status"] == "ACTIVE":
                    kb_id = kb["knowledgeBaseId"]
                    ds_resp = _get_bedrock_agent().list_data_sources(knowledgeBaseId=kb_id)
                    ds_list = ds_resp.get("dataSourceSummaries", [])
                    if ds_list:
                        return {"kb_id": kb_id, "data_source_id": ds_list[0]["dataSourceId"]}
    except Exception:
        log.debug("KB 조회 실패 (무시): kb_name=%s", kb_name, exc_info=True)
    return None


# ──────────────────────────────────────────────
# Ingestion 실행 / 폴링
# ──────────────────────────────────────────────
def start_ingestion(kb_id: str, data_source_id: str) -> str:
    """Ingestion Job 실행. 반환: ingestion_job_id."""
    resp = _get_bedrock_agent().start_ingestion_job(
        knowledgeBaseId=kb_id,
        dataSourceId=data_source_id,
    )
    return resp["ingestionJob"]["ingestionJobId"]


def is_ingestion_in_progress(kb_id: str, data_source_id: str) -> bool:
    """진행 중인 ingestion job 이 있는지 확인 (팀 KB 동시성 방지용)."""
    try:
        resp = _get_bedrock_agent().list_ingestion_jobs(
            knowledgeBaseId=kb_id,
            dataSourceId=data_source_id,
            sortBy={"attribute": "STARTED_AT", "order": "DESCENDING"},
            maxResults=1,
        )
        jobs = resp.get("ingestionJobSummaries", [])
        if jobs and jobs[0]["status"] in ("STARTING", "IN_PROGRESS"):
            return True
    except Exception:
        log.debug(
            "ingestion 상태 조회 실패 (무시): kb_id=%s ds_id=%s",
            kb_id, data_source_id, exc_info=True,
        )
    return False


# ──────────────────────────────────────────────
# 누적 문서 수 카운트 / 상한 검증
# ──────────────────────────────────────────────
# 변환본(원본은 originals/ 에 별도 보관) — 논리 문서 카운트에서 제외
_CONVERTED_SUFFIXES_COUNT = ("_pptx.json", "_xlsx.md", "_pdf.md")

# 같은 prefix(개인/팀 KB)의 업로드를 직렬화해 문서 수 상한 검증(read-modify-write)의
# TOCTOU 를 차단. 단일 백엔드 프로세스 기준 — 다중 프로세스 환경은 미보장.
_prefix_locks: dict[str, threading.Lock] = {}
_prefix_locks_guard = threading.Lock()


def _prefix_lock(prefix: str) -> threading.Lock:
    """prefix 별 락 인스턴스를 반환(없으면 생성)."""
    with _prefix_locks_guard:
        lock = _prefix_locks.get(prefix)
        if lock is None:
            lock = threading.Lock()
            _prefix_locks[prefix] = lock
        return lock


def _existing_logical_docs(bucket: str, raw_prefix: str) -> set[str]:
    """KB 의 논리 문서 파일명 집합.

    - originals/ 의 모든 파일 = 논리 문서 (pptx·Upstage pdf/xlsx·분할 csv/xlsx 의 원본)
    - raw/ 의 평문 파일 = 논리 문서 (pdf-local·docx·txt·md·json·html)
      변환본(_pptx.json/_xlsx.md/_pdf.md)·분할본(_partN)은 제외 — 원본을 originals 로 셈.

    분할본·멀티시트가 몇 개의 객체로 흩어지든 원본 1개로만 집계되므로, 시트 수만큼
    중복 카운트되거나(=과다 집계) 재업로드가 신규로 잡히는 문제를 함께 차단한다.
    """
    root = raw_prefix.rsplit("raw/", 1)[0]  # users/{emp}/ 또는 teams/{dept}/
    logical: set[str] = set()
    paginator = _get_s3().get_paginator("list_objects_v2")
    for page in paginator.paginate(Bucket=bucket, Prefix=root):
        for obj in page.get("Contents", []) or []:
            rel = obj["Key"][len(root):]          # raw/... 또는 originals/...
            sub, _, filename = rel.partition("/")
            if not filename:
                continue
            if sub == "originals":
                logical.add(filename)
            elif sub == "raw":
                if filename.endswith(_CONVERTED_SUFFIXES_COUNT) or is_tabular_part(filename):
                    continue
                logical.add(filename)
    return logical


def count_kb_docs(bucket: str, raw_prefix: str) -> int:
    """KB 의 논리적 문서 수(사용자가 올린 '파일' 수)."""
    return len(_existing_logical_docs(bucket, raw_prefix))


def _scope_from_prefix(prefix: str) -> Optional[str]:
    """업로드 prefix 로 KB scope 판별. 'users{env}/'→personal, 'teams{env}/'→team, 그 외 None."""
    if prefix.startswith(f"{kb_base('personal')}/"):
        return "personal"
    if prefix.startswith(f"{kb_base('team')}/"):
        return "team"
    return None  # shared 등은 상한 미적용


def enforce_kb_doc_limit(bucket: str, prefix: str, new_filenames: list[str]) -> None:
    """누적 문서 수 상한 검증. 초과 시 ValueError (배치당 개수 제한과 별개).

    이미 존재하는 동일 파일명(덮어쓰기)은 신규로 세지 않는다 — 한도가 찬 상태에서
    기존 문서 갱신이 잘못 거부되는 것을 방지. 상한 미정의 scope(shared 등)는 스킵.
    """
    scope = _scope_from_prefix(prefix)
    cap = KB_MAX_DOCS.get(scope) if scope else None
    if not cap:
        return
    existing = _existing_logical_docs(bucket, prefix)
    incoming = {Path(f).name for f in new_filenames}
    added = incoming - existing
    if len(existing) + len(added) > cap:
        label = {"personal": "개인", "team": "팀"}.get(scope, scope)
        raise ValueError(
            f"{label} 지식베이스에는 최대 {cap}개까지 업로드할 수 있습니다. "
            f"(현재 {len(existing)}개, 신규 {len(added)}개)"
        )


# ──────────────────────────────────────────────
# 파일 업로드 / 삭제 (S3)
# ──────────────────────────────────────────────
def _stash_original(
    bucket: str, prefix: str, file_bytes: bytes, filename: str, uploaded_uris: list[str],
) -> None:
    """변환 대상(pptx/pdf 등)의 원본을 raw/ 밖 originals/ 에 보관(Bedrock 인덱싱 제외).
    롤백 시 orphan 으로 남지 않도록 업로드 URI 목록에 기록."""
    originals_key = f"{get_originals_prefix(prefix)}{filename}"
    _get_s3().put_object(Bucket=bucket, Key=originals_key, Body=file_bytes)
    uploaded_uris.append(f"s3://{bucket}/{originals_key}")


def _validate_upload_files(files: list[tuple[bytes, str]]) -> None:
    """업로드 파일의 개수(≤5)·형식·크기 검증. S3 접근 없는 순수 검증이라
    staging 적재(stage_raw_files)와 색인 적재(upload_files_to_kb) 양쪽에서 재사용."""
    if len(files) > 5:
        raise ValueError(
            f"한 번에 최대 5개 파일만 업로드 가능합니다. (요청: {len(files)}개)"
        )
    for file_bytes, filename in files:
        ext = Path(filename).suffix.lower()
        if ext not in SUPPORTED_EXTENSIONS:
            raise ValueError(f"지원하지 않는 파일 형식: {filename}")
        validate_file_size(file_bytes, filename)


# PDF Upstage 변환 동시 실행 수. Upstage rate limit 을 고려한 보수적 상한.
_PDF_CONVERT_MAX_WORKERS = 3


def _convert_pdfs_parallel(files: list[tuple[bytes, str]]) -> dict[int, tuple[bytes, str]]:
    """업로드 목록 중 PDF(Upstage 모드)를 병렬 변환. 반환: {files 인덱스: (md_bytes, md_name)}.

    네트워크 바운드 Upstage 호출을 최대 _PDF_CONVERT_MAX_WORKERS 개 동시 실행해
    직렬 합산 지연을 줄인다(락·S3 쓰기 밖에서 수행). 변환 실패한 파일은 결과에서
    제외 → 호출자(upload_files_to_kb)가 원본 PDF 색인으로 폴백. files 인덱스로
    매핑해 동일 파일명 충돌을 피한다.
    """
    if not pdf_via_upstage_enabled():
        return {}
    targets = [
        (i, fb, fn) for i, (fb, fn) in enumerate(files)
        if Path(fn).suffix.lower() == ".pdf"
    ]
    if not targets:
        return {}

    result: dict[int, tuple[bytes, str]] = {}
    with ThreadPoolExecutor(max_workers=min(_PDF_CONVERT_MAX_WORKERS, len(targets))) as ex:
        fut_to_idx = {
            ex.submit(convert_pdf_to_markdown, fb, fn): i for i, fb, fn in targets
        }
        for fut in as_completed(fut_to_idx):
            i = fut_to_idx[fut]
            try:
                result[i] = fut.result()
            except Exception:
                log.warning(
                    "Upstage PDF 변환 실패, 원본 PDF 색인으로 폴백: %s",
                    files[i][1], exc_info=True,
                )
    return result


def upload_files_to_kb(
    bucket: str,
    prefix: str,
    files: list[tuple[bytes, str]],
    with_rollback: bool = False,
) -> list[str]:
    """파일 목록을 S3 {bucket}/{prefix} 에 업로드. 최대 5개.

    - pptx: 원본은 originals/ 로, JSON 변환본은 raw/ 로 업로드
    - xlsx/csv: ROWS_PER_SPLIT 단위로 분할 업로드
    - 그 외: 단일 업로드 (형식별 크기 제한 검증)
    - with_rollback=True 시 도중 실패하면 부분 업로드분을 모두 삭제

    반환: 업로드된 S3 URI 목록 (pptx 의 originals/ 원본 URI 포함).
    """
    # 형식/크기 검증은 S3 접근이 없어 락 밖에서 선수행.
    _validate_upload_files(files)

    # PDF Upstage 변환을 락·S3 쓰기 밖에서 미리 병렬 수행(직렬 합산 지연 단축).
    # 변환 실패분은 pdf_md 에 없어 아래 루프가 원본 PDF 색인으로 폴백.
    pdf_md = _convert_pdfs_parallel(files)

    # 상한 검증(count)→S3 쓰기를 prefix 락으로 직렬화해 TOCTOU(동시 업로드가 같은
    # count 를 읽고 둘 다 통과)를 차단. 단일 백엔드 프로세스 기준.
    with _prefix_lock(prefix):
        # 누적 문서 수 상한 (개인/팀). 초과 시 ValueError.
        enforce_kb_doc_limit(bucket, prefix, [fn for _, fn in files])

        uploaded_uris: list[str] = []
        try:
            for idx, (file_bytes, filename) in enumerate(files):
                ext = Path(filename).suffix.lower()

                if ext in CONVERTIBLE_EXTS:
                    # 원본은 originals/ 에 보관(인덱싱 제외), 변환본(json)만 raw/ 에 색인.
                    _stash_original(bucket, prefix, file_bytes, filename, uploaded_uris)
                    file_bytes, filename = convert_pptx_to_json(file_bytes, filename)
                    ext = ".json"
                elif ext == ".pdf" and pdf_via_upstage_enabled():
                    # PDF 는 위에서 병렬 변환됨(pdf_md). 원본 PDF 는 originals/ 에 보관,
                    # 변환본 _pdf.md 만 raw/ 에 색인. 변환 실패분(pdf_md 에 없음)은 원본
                    # PDF 를 그대로 색인 → Lambda parse_pdf 폴백.
                    md = pdf_md.get(idx)
                    if md is not None:
                        _stash_original(bucket, prefix, file_bytes, filename, uploaded_uris)
                        file_bytes, filename = md
                        ext = ".md"

                # xlsx 는 여기(개인/팀)서 Upstage 변환하지 않고 pandas 분할(TABULAR_EXTS)로 처리.
                # xlsx→Upstage 는 공용 KB CLI(shared_kb_manager)에서만 적용하는 정책 — 의도된 비대칭.
                if ext in TABULAR_EXTS:
                    # 분할본(_partN)과 별개로 원본을 originals/ 에 보관 → 논리 문서 카운트·
                    # 목록·삭제가 멀티시트/분할 여부와 무관하게 '파일 1개'로 일관 처리.
                    _stash_original(bucket, prefix, file_bytes, filename, uploaded_uris)
                    uris = split_and_upload_tabular(bucket, prefix, file_bytes, filename)
                else:
                    key = f"{prefix}{filename}"
                    _get_s3().put_object(Bucket=bucket, Key=key, Body=file_bytes)
                    uris = [f"s3://{bucket}/{key}"]
                uploaded_uris.extend(uris)

        except Exception:
            if with_rollback and uploaded_uris:
                log.warning("S3 업로드 실패, 롤백 시작: %d개 삭제", len(uploaded_uris))
                for uri in uploaded_uris:
                    key = uri.replace(f"s3://{bucket}/", "")
                    try:
                        _get_s3().delete_object(Bucket=bucket, Key=key)
                    except Exception as del_err:
                        log.warning("S3 롤백 실패: %s, %s", key, del_err)
            raise

    return uploaded_uris


def _delete_keys_quietly(bucket: str, keys: list[str]) -> None:
    """주어진 S3 키들을 best-effort 삭제(실패는 경고만). staging 정리·롤백 공용."""
    for key in keys:
        try:
            _get_s3().delete_object(Bucket=bucket, Key=key)
        except Exception as e:
            log.warning("S3 삭제 실패(무시): %s, %s", key, e)


def stage_raw_files(
    bucket: str,
    prefix: str,
    files: list[tuple[bytes, str]],
) -> list[str]:
    """원본 바이트를 staging/ 에만 적재(변환·색인 X) — 업로드 HTTP 요청용.

    변환(Upstage 등)을 요청 밖(백그라운드)으로 미뤄 다중 PDF 등에서 프록시
    타임아웃(504)을 방지. staging 적재 *전에* 형식/크기·누적 상한을 선검증해,
    상한 초과 시 S3 에 아무것도 올리지 않고 즉시 거부(고아 방지). 부분 적재 실패
    시 이미 올린 분을 롤백.

    prefix: raw/ prefix (적재 위치는 그 형제 staging/).
    반환: staging 에 적재된 파일명 목록.
    """
    _validate_upload_files(files)
    # 누적 문서 수 상한을 staging 전에 선검증(빠른 거부, 사용자 즉시 피드백).
    # 최종 보장은 색인 시 upload_files_to_kb 가 prefix 락 하에서 재검증(동시 race 방어).
    enforce_kb_doc_limit(bucket, prefix, [fn for _, fn in files])

    staging = get_staging_prefix(prefix)
    staged_names: list[str] = []
    try:
        for file_bytes, filename in files:
            _get_s3().put_object(Bucket=bucket, Key=f"{staging}{filename}", Body=file_bytes)
            staged_names.append(filename)
    except Exception:
        _delete_keys_quietly(bucket, [f"{staging}{n}" for n in staged_names])
        raise
    return staged_names


def process_staged_files(
    bucket: str,
    prefix: str,
    names: list[str],
) -> list[str]:
    """staging/ 의 원본을 읽어 변환+raw/+originals/ 적재(색인 준비) 후 staging 정리.

    백그라운드(on_upload_complete)에서 호출 — Upstage 변환 등 무거운 작업이 여기서
    일어난다(HTTP 요청 밖). 변환·적재는 upload_files_to_kb 를 그대로 재사용하며
    with_rollback=True 라 실패 시 부분 raw/originals 가 정리된다. staging/ 원본은
    성공·실패와 무관하게 정리해 고아를 남기지 않는다.

    prefix: raw/ prefix. names: staging 에 적재된 원본 파일명(=업로드 파일명).
    반환: 색인된 raw/originals URI 목록.
    """
    staging = get_staging_prefix(prefix)
    staged_keys = [f"{staging}{name}" for name in names]
    try:
        files: list[tuple[bytes, str]] = []
        for key, name in zip(staged_keys, names):
            obj = _get_s3().get_object(Bucket=bucket, Key=key)
            files.append((obj["Body"].read(), name))
        return upload_files_to_kb(bucket, prefix, files, with_rollback=True)
    finally:
        # 성공: 원본 staging 정리 / 실패: upload_files_to_kb 가 raw·originals 롤백,
        # 여기서 staging 정리 → 어느 경로든 고아 없음.
        _delete_keys_quietly(bucket, staged_keys)


def delete_files_from_kb(bucket: str, prefix: str, filenames: list[str]) -> None:
    """선택된 파일들을 S3 에서 삭제.

    pptx 의 경우 원본(originals/) 과 인덱싱본(raw/_pptx.json) 둘 다 삭제.
    삭제 후 ingestion job 을 실행해야 Bedrock 이 변경을 감지하여
    S3 Vectors 의 해당 파일 벡터를 제거 (호출자가 직접 트리거).

    xlsx/csv 분할본(_partN)은 개수·시트 구성을 미리 알 수 없으므로 raw/ 를 나열해
    originals/ 기준 소유 판정으로 해당 파일 소유분만 삭제 (prefix 충돌 파일 보존).

    S3 delete_object 는 멱등이므로 키가 없어도 예외를 던지지 않음.
    """
    originals = get_originals_prefix(prefix)
    original_stems = _list_original_stems(bucket, originals)
    keys_to_delete: list[str] = []
    for filename in filenames:
        ext = Path(filename).suffix.lower()
        stem = Path(filename).stem
        # 원본 보관본(pptx/Upstage pdf·xlsx/분할 csv·xlsx). 평문 파일엔 없어도 멱등 no-op.
        keys_to_delete.append(f"{originals}{filename}")
        if ext == ".pptx":
            keys_to_delete.append(f"{prefix}{stem}_pptx.json")
        elif ext == ".xlsx":
            keys_to_delete.append(f"{prefix}{stem}_xlsx.md")          # Upstage 변환본
            keys_to_delete.extend(                                    # local 분할본(_partN)
                _owned_part_keys(bucket, prefix, stem, ".xlsx", original_stems)
            )
        elif ext == ".csv":
            keys_to_delete.extend(                                    # csv 분할본(_partN)
                _owned_part_keys(bucket, prefix, stem, ".csv", original_stems)
            )
        elif ext == ".pdf":
            keys_to_delete.append(f"{prefix}{stem}_pdf.md")           # Upstage 변환본
            keys_to_delete.append(f"{prefix}{filename}")             # 미변환 원본 색인본
        else:
            keys_to_delete.append(f"{prefix}{filename}")

    for key in dict.fromkeys(keys_to_delete):  # 중복 키 제거(순서 보존)
        _get_s3().delete_object(Bucket=bucket, Key=key)


# ──────────────────────────────────────────────
# Ingestion 상태 폴링
# ──────────────────────────────────────────────

# Bedrock failureReasons 는 실패 문서 목록을 통째로 담은 장문 사유를 반환할 때가
# 있어, 그대로 로그에 남기면 JSON 로그 1줄이 수십 KB 로 비대해진다(파일 회전 가속·
# 대시보드 노이즈). 사유당·전체 길이를 제한한다.
_FAIL_REASON_MAX_CHARS = 300
_FAIL_DETAIL_MAX_CHARS = 800


def _cap_fail_detail(reasons: list) -> str:
    """failureReasons 를 사유당/전체 길이 제한으로 요약."""
    detail = "; ".join(str(r)[:_FAIL_REASON_MAX_CHARS] for r in (reasons or [])[:3])
    if len(detail) > _FAIL_DETAIL_MAX_CHARS:
        detail = detail[:_FAIL_DETAIL_MAX_CHARS] + f"…(+{len(detail) - _FAIL_DETAIL_MAX_CHARS}자 생략)"
    return detail


def poll_ingestion_status(
    kb_id: str,
    data_source_id: str,
    job_id: str,
    poll_interval: int | None = None,
    poll_timeout: int | None = None,
) -> str:
    """
    Ingestion Job 완료 여부 폴링.
    반환: 최종 상태 문자열.
    실패 시 상태에 사유를 포함: "FAILED: 사유..."
    """
    cfg = _cfg()
    if poll_interval is None:
        poll_interval = cfg.get("ingest_poll_interval", 5)
    if poll_timeout is None:
        poll_timeout = cfg.get("ingest_poll_timeout", 300)
    start = time.time()
    while time.time() - start < poll_timeout:
        resp = _get_bedrock_agent().get_ingestion_job(
            knowledgeBaseId=kb_id,
            dataSourceId=data_source_id,
            ingestionJobId=job_id,
        )
        job = resp["ingestionJob"]
        status = job["status"]
        stats = job.get("statistics", {})

        if status in ("COMPLETE", "FAILED", "STOPPED"):
            if status != "COMPLETE":
                reasons = job.get("failureReasons", [])
                if reasons:
                    detail = _cap_fail_detail(reasons)
                    log.warning("Bedrock Ingestion 실패: %s", detail)
                    return f"{status}: {detail}"
            num_failed = stats.get("numberOfDocumentsFailed", 0)
            if status == "COMPLETE" and num_failed > 0:
                num_indexed = stats.get("numberOfNewDocumentsIndexed", 0)
                num_scanned = stats.get("numberOfDocumentsScanned", 0)
                log.warning(
                    "Bedrock 부분 실패: %d개 문서 실패 (scanned=%d, indexed=%d)",
                    num_failed, num_scanned, num_indexed,
                )
                return f"COMPLETE_WITH_ERRORS: {num_indexed}개 성공, {num_failed}개 실패"
            return status
        time.sleep(poll_interval)
    raise TimeoutError(f"Ingestion 완료 대기 타임아웃: job_id={job_id}")
