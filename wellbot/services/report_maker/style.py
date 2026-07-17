"""문서 스타일 학습 — PPTX/PDF/이미지 파싱 + LLM 스타일 분석.

legacy storage.py(문서 파싱) + memory.py(스타일 분석/서술)에서 순수 파싱 로직은
원문 보존, LLM 호출부만 Converse 헬퍼(bedrock)로 재작성했다.
"""

from __future__ import annotations

import logging
from pathlib import Path

from pptx import Presentation

from wellbot.services.report_maker import bedrock
from wellbot.services.report_maker.config import DOC_EXTS, IMAGE_EXTS, get_config

log = logging.getLogger(__name__)

_MEDIA_TYPE_MAP = {
    ".jpg": "jpeg", ".jpeg": "jpeg", ".png": "png", ".gif": "gif", ".webp": "webp",
}


def _sanitize(text: str) -> str:
    """서로게이트 문자 제거(JSON/S3 안전)."""
    return "".join(c for c in str(text) if not (0xD800 <= ord(c) <= 0xDFFF))


def is_image_file(path: str) -> bool:
    return Path(path).suffix.lower() in IMAGE_EXTS


def is_doc_file(path: str) -> bool:
    return Path(path).suffix.lower() in DOC_EXTS


# ══════════════════════════════════════════════════════════════
# 이미지 텍스트 추출 (Bedrock Converse vision) — Converse 포맷으로 재작성
# ══════════════════════════════════════════════════════════════
def extract_text_from_image(image_path: str) -> str:
    ext = Path(image_path).suffix.lower()
    fmt = _MEDIA_TYPE_MAP.get(ext)
    if not fmt:
        raise ValueError(f"지원하지 않는 이미지 형식: {ext}")
    with open(image_path, "rb") as f:
        data = f.read()

    # 공유 Converse 헬퍼 사용 — ThrottlingException 재시도·에러처리 일원화(실패 시 "").
    extracted = bedrock.call_vision(
        data,
        fmt,
        (
            "이 이미지에서 모든 텍스트를 추출해주세요. 슬라이드 제목, 본문, "
            "bullet 항목, 표 내용 등 보이는 텍스트를 원본 구조 그대로 출력하세요. "
            "설명이나 부연은 하지 말고 텍스트만 출력하세요."
        ),
        get_config().max_tokens_vision,
    )
    log.info("이미지 텍스트 추출 완료 (%d자)", len(extracted))
    return extracted


# ══════════════════════════════════════════════════════════════
# 순수 파싱 (PPTX/PDF) — legacy storage.py 원문 보존
# ══════════════════════════════════════════════════════════════
def _zone_by_ratio(cx: float, width: float) -> str:
    if not width or width <= 0:
        return "중"
    r = cx / width
    return "좌" if r < 0.4 else ("우" if r > 0.6 else "중")
    
# ──────────────────────────────────────────────────────────────
# PPTX 파싱
# ──────────────────────────────────────────────────────────────
def _extract_table_content(shape) -> dict:
    rows = [[cell.text.strip() for cell in row.cells] for row in shape.table.rows]
    rows = [r for r in rows if any(r)]
    if not rows:
        return {}
    return {
        "headers":   rows[0],
        "rows":      rows[1:],
        "row_count": len(rows) - 1,
        "col_count": len(rows[0]),
        "structure": (
            "비교형" if len(rows[0]) == 2 else
            "다중비교형" if len(rows[0]) >= 3 and len(rows) <= 5 else
            "목록형"
        ),
    }


def _find_slide_title(slide) -> str:
    """슬라이드 제목을 견고하게 찾는다.
    1순위: 제목 플레이스홀더(TITLE=1, CENTER_TITLE=0)
    2순위: 그 슬라이드에서 '가장 큰 폰트'의 텍스트 (절대 크기 아님 — 상대 최대).
           단 본문(그보다 작은 폰트)이 존재해 계층이 있을 때만. 크기 정보가
           전혀 없거나 모두 같은 크기면 제목을 특정하지 않는다(오인 방지)."""
    # 1) 제목 플레이스홀더
    for shape in slide.shapes:
        if shape.has_text_frame and shape.is_placeholder:
            try:
                if int(shape.placeholder_format.type) in (0, 1):
                    t = shape.text_frame.text.strip()
                    if t:
                        return t
            except Exception:
                pass
    # 2) 폴백: 문단별 (텍스트, 폰트pt) 수집
    items = []  # (text, size_pt)
    for shape in slide.shapes:
        if not shape.has_text_frame:
            continue
        for para in shape.text_frame.paragraphs:
            t = para.text.strip()
            if not t:
                continue
            sz = 0
            try:
                if para.runs and para.runs[0].font.size:
                    sz = para.runs[0].font.size.pt
                elif para.font.size:
                    sz = para.font.size.pt
            except Exception:
                sz = 0
            items.append((t, sz))
 
    sized = [(t, s) for t, s in items if s and s > 0]
    if not sized:
        return ""                      # 크기 정보 없음 → 제목 특정 불가
    max_sz = max(s for _, s in sized)
    body_szs = [s for _, s in sized if s < max_sz]
    if not body_szs:
        return ""                      # 전부 같은 크기 → 제목 계층 없음(오인 방지)
    # 가장 큰 폰트의 첫 텍스트를 제목으로 (문서마다 크기 달라도 상대적으로 최대면 제목)
    for t, s in sized:
        if s == max_sz:
            return t
    return ""
 
 
def _extract_pptx(pptx_path: str) -> dict:
    prs, slides_info = Presentation(pptx_path), []
    sw = prs.slide_width or 1
    for i, slide in enumerate(prs.slides, 1):
        sd = {"slide_num": i, "title": "", "bullets": [], "paragraphs": [],
              "tables": [], "content_items": []}
 
        # ── 제목 먼저 확정 (플레이스홀더 우선 → 폰트 폴백) ──
        sd["title"] = _find_slide_title(slide)
        _title_used = sd["title"]   # 본문에서 제외할 제목 텍스트
 
        for shape in slide.shapes:
            # 도형 x중심 → 좌/중/우
            try:
                cx = shape.left + (shape.width or 0) / 2 if shape.left is not None else None
                zone = _zone_by_ratio(cx, sw) if cx is not None else "중"
            except Exception:
                zone = "중"
 
            if shape.has_table:
                t = _extract_table_content(shape)
                if t:
                    sd["tables"].append(t)
                    sd["content_items"].append({"zone": zone, "text": "[표]"})
            elif shape.has_text_frame:
                for para in shape.text_frame.paragraphs:
                    text = para.text.strip()
                    if not text:
                        continue
                    if text == _title_used:      # 이미 제목으로 쓴 텍스트는 본문에서 제외
                        continue
                    if any(text.startswith(c) for c in ["•","-","▪","○","·"]) or para.level > 0:
                        sd["bullets"].append(text)
                    else:
                        sd["paragraphs"].append(text)
                    sd["content_items"].append({"zone": zone, "text": text})
        slides_info.append(sd)
        print(f"[제목확인] 슬라이드{i}: title='{sd['title']}' | 본문 {len(sd['content_items'])}개")
 
    bullet_cnt = sum(len(s["bullets"])    for s in slides_info)
    para_cnt   = sum(len(s["paragraphs"]) for s in slides_info)
    all_tables = [t for s in slides_info for t in s["tables"]]
    table_pats: dict = {}
    for t in all_tables:
        table_pats[t["structure"]] = table_pats.get(t["structure"], 0) + 1
 
    return {
        "slide_details": [
            {"slide_num": s["slide_num"], "title": s["title"],
             "content": s["bullets"] + s["paragraphs"], "tables": s["tables"],
             "content_items": s["content_items"]}
            for s in slides_info
            if s["title"] or s["bullets"] or s["paragraphs"] or s["tables"]
        ],
        "toc":            [s["title"] for s in slides_info if s["title"]],
        "slide_count":    len(prs.slides),
        "writing_style":  "개조식" if bullet_cnt > para_cnt else "서술형",
        "table_count":    len(all_tables),
        "table_patterns": table_pats,
    }
# ──────────────────────────────────────────────────────────────
# PDF 파싱 (pdfplumber)
# ──────────────────────────────────────────────────────────────

def _layout_from_spans(spans: list, page_w: float) -> str:
    if not spans:
        return "title_only"
    # 판정에 쓸 박스: 너무 넓은 것(전폭·배경)과 너무 좁은 것(제목·라벨·아이콘)은 제외
    core = [s for s in spans
            if page_w * 0.15 <= (s[1] - s[0]) < page_w * 0.70]
    n = len(core)
    if n == 0:
        return "single"
    band_lo, band_hi = page_w * 0.45, page_w * 0.55
    crossing   = [s for s in core if s[0] < band_hi and s[1] > band_lo]
    left_only  = any(r <= band_lo for l, r in core)
    right_only = any(l >= band_hi for l, r in core)
    if left_only and right_only and len(crossing) <= max(int(n * 0.25), 1):
        return "two_col"
    return "single"

def _extract_layout_from_page(page) -> dict:
    pw = page.width
    words = page.extract_words(extra_attrs=["size"]) or []
    font_sizes = [w.get("size", 0) for w in words if w.get("size")]
    title_size = max(font_sizes) if font_sizes else 0
 
    title = ""
    content_items = []
    for w in words:
        text = w["text"].strip()
        if not text:
            continue
        # 제목(가장 큰 글씨 첫 등장) 분리
        if w.get("size", 0) >= title_size * 0.95 and not title:
            title = text
            continue
        cx = (w["x0"] + w["x1"]) / 2
        content_items.append({"zone": _zone_by_ratio(cx, pw), "text": text})
 
    return {"title": title, "content_items": content_items}

def _extract_pdf(pdf_path: str) -> dict:
    import pdfplumber
    pages_info: list = []
 
    with pdfplumber.open(pdf_path) as pdf:
        for i, page in enumerate(pdf.pages, 1):
            pw = page.width or 1
            pd = {"slide_num": i, "title": "", "content_items": [], "tables": []}
 
            # ── 표 추출 (기존 유지) + 표 위치 태깅 ──
            table_bboxes = []
            try:
                found = page.find_tables()
                table_bboxes = [t.bbox for t in found]  # (x0,y0,x1,y1)
            except Exception:
                table_bboxes = []
 
            for ti, table in enumerate(page.extract_tables()):
                if not table:
                    continue
                rows = [[cell.strip() if cell else "" for cell in row] for row in table]
                rows = [r for r in rows if any(r)]
                if not rows:
                    continue
                pd["tables"].append({
                    "headers":   rows[0], "rows": rows[1:],
                    "row_count": len(rows) - 1, "col_count": len(rows[0]),
                    "structure": (
                        "비교형"     if len(rows[0]) == 2 else
                        "다중비교형"  if len(rows[0]) >= 3 and len(rows) <= 5 else
                        "목록형"
                    ),
                })
                # 표의 x중심으로 zone 태깅 (bbox가 있으면)
                if ti < len(table_bboxes):
                    bx0, _, bx1, _ = table_bboxes[ti]
                    cx = (bx0 + bx1) / 2
                    z = "좌" if cx / pw < 0.4 else ("우" if cx / pw > 0.6 else "중")
                else:
                    z = "중"
                pd["content_items"].append({"zone": z, "text": "[표]"})
 
            # ── 본문: 새 _extract_layout_from_page 에서 title + content_items(zone) ──
            layout = _extract_layout_from_page(page)
            pd["title"] = layout.get("title", "")
            pd["content_items"] += layout.get("content_items", [])
 
            pages_info.append(pd)
 
    # writing_style 계산용: bullet vs paragraph (content_items 텍스트에서 추정)
    def _is_bullet(t):
        return any(t.startswith(c) for c in ["•","-","▪","○","·","–"])
    bullet_cnt = sum(1 for p in pages_info for it in p["content_items"]
                     if it["text"] != "[표]" and _is_bullet(it["text"]))
    para_cnt   = sum(1 for p in pages_info for it in p["content_items"]
                     if it["text"] != "[표]" and not _is_bullet(it["text"]))
 
    all_tables = [t for p in pages_info for t in p["tables"]]
    table_pats: dict = {}
    for t in all_tables:
        table_pats[t["structure"]] = table_pats.get(t["structure"], 0) + 1
 
    return {
        "slide_details": [
            {"slide_num": p["slide_num"], "title": p["title"],
             "content": [it["text"] for it in p["content_items"] if it["text"] != "[표]"],
             "tables": p["tables"],
             "content_items": p["content_items"]}         # ← LLM 판정용 위치정보
            for p in pages_info
            if p["title"] or p["content_items"] or p["tables"]
        ],
        "toc":            [p["title"] for p in pages_info if p["title"]],
        "slide_count":    len(pages_info),
        "writing_style":  "개조식" if bullet_cnt > para_cnt else "서술형",
        "table_count":    len(all_tables),
        "table_patterns": table_pats,

    }


def extract_doc_style(file_path: str) -> dict:
    ext = Path(file_path).suffix.lower()
    if ext == ".pdf":
        return _extract_pdf(file_path)
    elif ext == ".pptx":
        return _extract_pptx(file_path)
    else:
        raise ValueError(f"지원하지 않는 파일 형식: {ext}")
           


# ══════════════════════════════════════════════════════════════
# 스타일 분석 (LLM) — 프롬프트 원문 보존, 호출부 Converse 로 재작성
# ══════════════════════════════════════════════════════════════
def analyze_style_with_claude(doc_style: dict) -> dict:
    toc_text = "\n".join(f"{i+1}. {t}" for i, t in enumerate(doc_style["toc"]))
    slide_count = doc_style.get("slide_count", len(doc_style.get("slide_details", [])))

    slides_text = ""
    for s in doc_style["slide_details"]:
        slides_text += f"\n[슬라이드 {s['slide_num']}] {s['title']}\n"
        items = s.get("content_items", [])
        if items:
            zones = {"좌": [], "중": [], "우": []}
            for it in items:
                zones.get(it.get("zone", "중"), zones["중"]).append(it.get("text", ""))
            for z in ("좌", "중", "우"):
                if zones[z]:
                    joined = " ".join(t for t in zones[z] if t)
                    slides_text += f"  ({z}) {joined[:300]}\n"
        else:
            for txt in s.get("content", []):
                slides_text += f"  - {txt}\n"
        for t in s.get("tables", []):
            slides_text += f"  [표] {t['structure']} | 헤더: {t['headers']}\n"
            for row in t["rows"][:3]:
                slides_text += f"      {row}\n"

    prompt = f"""다음 PPT/PDF 문서를 분석해서 작성 스타일을 JSON으로만 답변하세요.

[슬라이드 수]
{slide_count}장

[슬라이드 제목 시퀀스]
{toc_text}

[슬라이드별 세부 내용]
{slides_text}

[레이아웃 판정 지침]
각 슬라이드가 좌우 2단(two_col)인지 단일 영역(single)인지 판정하세요.
- (좌)와 (우)에 서로 다른 성격의 내용이 나뉘어 있으면 two_col
- 한 흐름으로 이어지거나, 한쪽(또는 중앙)에만 내용이 있으면 single
- 좌우로 갈렸어도 단순히 같은 목록이 이어진 것뿐이면 single
- 표 하나로 페이지가 채워지면 single

JSON 형식:
{{
  "문서_작성_스타일": {{
    "문서유형": "...",
    "문서목적": "...",
    "영역별_bullet_수": "...",
    "불렛_구조": "... (예: 1단계만 사용 / 2단계: 주장+근거 / 3단계: 주장+근거+사례)",
    "전개_방식": "..."
  }},
  "표현규칙": {{
    "문장종결": "...",
    "보고서_톤": "...",
    "약어_사용_방식": "...",
    "선호_문체": "...",
    "근거_제시_방식": "..."
  }},
  "출력형식": {{
    "슬라이드_분량": "{slide_count}장",
    "표_활용": "..."
  }},
  "페이지별_레이아웃": [
    {{"슬라이드": 1, "레이아웃": "single 또는 two_col"}}
  ]
}}"""

    result = bedrock.call_json(prompt, get_config().max_tokens_style)

    per_page = result.get("페이지별_레이아웃", [])
    counts: dict = {}
    for p in per_page:
        lt = p.get("레이아웃", "single")
        lt = lt if lt in ("single", "two_col") else "single"
        counts[lt] = counts.get(lt, 0) + 1
    total = sum(counts.values()) or 1
    layout_ratio = {k: f"{round(v/total*100)}%" for k, v in counts.items()}
    doc_style["layout_ratio"] = layout_ratio
    result.setdefault("출력형식", {})["선호_슬라이드_구성"] = (
        ", ".join(f"{k} {v}" for k, v in layout_ratio.items()) if layout_ratio else "N/A"
    )
    return result


def build_style_desc(doc_style: dict, analysis: dict) -> str:
    doc  = analysis.get("문서_작성_스타일", {})
    expr = analysis.get("표현규칙", {})
    fmt  = analysis.get("출력형식", {})

    lr = doc_style.get("layout_ratio", {})
    if lr:
        merged = {}
        for k, v in lr.items():
            key = "single" if k in ("single", "title_only") else k
            pct = int(str(v).rstrip("%") or 0)
            merged[key] = merged.get(key, 0) + pct
        order = sorted(merged.items(), key=lambda x: -x[1])
        primary = order[0][0]
        ratio_str = ", ".join(f"{k} {v}%" for k, v in order)
        layout_hint = (
            f"- 레이아웃: 주로 {primary}로 구성하되, 내용·스토리라인상 다른 레이아웃이 "
            f"자연스러운 페이지는 그에 맞게 둔다. 억지로 비율을 맞추지 말고 페이지 내용에 따라 "
            f"판단한다. (참고 비율: {ratio_str})\n"
        )
    else:
        layout_hint = ""

    return (
        "[문서 작성 스타일]\n"
        f"* 문서유형: {doc.get('문서유형', 'N/A')}\n"
        f"* 문서목적: {doc.get('문서목적', 'N/A')}\n"
        f"* 영역별 bullet 수: {doc.get('영역별_bullet_수', 'N/A')}\n"
        f"* 전개 방식: {doc.get('전개_방식', 'N/A')}\n\n"
        "[표현규칙]\n"
        f"* 문장종결: {expr.get('문장종결', 'N/A')}\n"
        f"* 보고서 톤: {expr.get('보고서_톤', 'N/A')}\n"
        f"* 약어 사용 방식: {expr.get('약어_사용_방식', 'N/A')}\n"
        f"* 선호 문체: {expr.get('선호_문체', 'N/A')}\n"
        f"* 근거 제시 방식: {expr.get('근거_제시_방식', 'N/A')}\n\n"
        "[출력형식]\n"
        f"* 슬라이드 분량: {fmt.get('슬라이드_분량', 'N/A')}\n"
        f"* 선호 슬라이드 구성: {fmt.get('선호_슬라이드_구성', 'N/A')}\n"
        f"{layout_hint}" 
        f"* 표 활용: {fmt.get('표_활용', 'N/A')}"
    )
