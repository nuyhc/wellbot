# storage.py — S3 파일 관리 + PPTX/PDF/이미지 파싱 + AgentCore 저장
import json, tempfile, base64
import uuid
import re
from pathlib import Path
from datetime import datetime, timezone, timedelta
from pptx import Presentation
from config import KST, S3_BUCKET, S3_PREFIX, MEMORY_ID, MODEL_ID, bedrock, ac_client, s3_client

IMAGE_EXTS = {".jpg", ".jpeg", ".png", ".gif", ".webp"}
DOC_EXTS = {".pptx", ".pdf"}
newline ='\n'

# ──────────────────────────────────────────────────────────────
# 유틸
# ──────────────────────────────────────────────────────────────
def _sanitize(text: str) -> str:
    return "".join(c for c in str(text) if not (0xD800 <= ord(c) <= 0xDFFF))


# ──────────────────────────────────────────────────────────────
# S3 파일 관리
# ──────────────────────────────────────────────────────────────
def upload_to_s3(local_path: str, user_id: str, template_id: str, folder: str = "input") -> str:
    ts = datetime.now(tz=KST).strftime("%y%m%d%H%M%S")
    src_path = Path(local_path)
    safe_name = src_path.name.encode("utf-8", errors="ignore").decode("utf-8")
    key = f"{S3_PREFIX}/{user_id}/{template_id}/{folder}/{ts}_{safe_name}"
    s3_client.upload_file(str(src_path), S3_BUCKET, key)
    print(f"   S3 업로드 완료: {key}")
    return key

def list_conversations(user_id: str, template_id: str) -> list:
    prefix = f"{S3_PREFIX}/{user_id}/{template_id}/conversations/"
    try:
        resp  = s3_client.list_objects_v2(Bucket=S3_BUCKET, Prefix=prefix)
        items = []
        for obj in resp.get("Contents", []):
            try:
                o    = s3_client.get_object(Bucket=S3_BUCKET, Key=obj["Key"])
                data = json.loads(o["Body"].read())
                items.append({
                    "session_id": data.get("session_id", ""),
                    "title":      data.get("title", "제목 없음"),
                    "saved_at":   data.get("saved_at", ""),
                })
            except Exception:
                continue
        return sorted(items, key=lambda x: x["saved_at"], reverse=True)
    except Exception as e:
        print(f"    대화 목록 조회 실패: {e}")
        return []

def upload_to_s3_style(local_path: str, user_id: str, template_id: str) -> str:
    ts        = datetime.now(tz=KST).strftime("%y%m%d%H%M%S")
    src_path  = Path(local_path)
    safe_name = src_path.name.encode("utf-8", errors="ignore").decode("utf-8")
    key       = f"{S3_PREFIX}/{user_id}/{template_id}/input/style_docs/{ts}_{safe_name}"
    s3_client.upload_file(str(src_path), S3_BUCKET, key)
    print(f"  S3 style_docs 업로드 완료: {key}")
    return key

def list_files_from_s3(user_id: str, template_id: str) -> list:
    prefix = f"{S3_PREFIX}/{user_id}/{template_id}/input/style_docs/"
    resp = s3_client.list_objects_v2(Bucket=S3_BUCKET, Prefix=prefix)
    files = [
        o for o in resp.get("Contents", [])
        if Path(o["Key"]).suffix.lower() in DOC_EXTS | IMAGE_EXTS
    ]
    return sorted(files, key=lambda o: o["LastModified"], reverse=True)

def download_file_from_s3(s3_key: str) -> str:
    suffix = Path(s3_key).suffix.lower()
    tmp = tempfile.NamedTemporaryFile(suffix=suffix, delete=False)
    s3_client.download_fileobj(S3_BUCKET, s3_key, tmp)
    tmp.flush()
    return tmp.name

def get_analyzed_history(user_id: str, template_id: str) -> set:
    key = f"{S3_PREFIX}/{user_id}/{template_id}/meta/analyzed.json"
    try:
        obj = s3_client.get_object(Bucket=S3_BUCKET, Key=key)
        return set(json.loads(obj["Body"].read()).get("analyzed", []))
    except s3_client.exceptions.NoSuchKey:
        return set()
    except Exception as e:
        print(f"    분석 이력 로드 실패: {e}")
        return set()

def save_analyzed_history(user_id: str, template_id: str, analyzed: set):
    key = f"{S3_PREFIX}/{user_id}/{template_id}/meta/analyzed.json"
    body = json.dumps({"analyzed": sorted(analyzed)}, ensure_ascii=False)
    s3_client.put_object(Bucket=S3_BUCKET, Key=key, Body=body,
                         ContentType="application/json")
    print(f"   분석 이력 저장 완료 ({len(analyzed)}개)")


def load_combined_style(user_id: str, template_id: str) -> str:
    key = f"{S3_PREFIX}/{user_id}/{template_id}/meta/combined_style.json"
    try:
        obj = s3_client.get_object(Bucket=S3_BUCKET, Key=key)
        text = json.loads(obj["Body"].read()).get("style_desc", "")
        if text:
            print("   combined_style.json 로드 완료")
        return text
    except s3_client.exceptions.NoSuchKey:
        return ""
    except Exception as e:
        print(f"    combined_style.json 로드 실패: {e}")
        return ""


def save_combined_style(user_id: str, template_id: str, new_style_desc: str):
    key = f"{S3_PREFIX}/{user_id}/{template_id}/meta/combined_style.json"
    
    # 기존 스타일 로드
    existing = load_combined_style(user_id, template_id)
    
    if not existing:
        # 최초 저장은 그냥 저장
        body = json.dumps({"style_desc": new_style_desc}, ensure_ascii=False)
        s3_client.put_object(Bucket=S3_BUCKET, Key=key, Body=body,
                             ContentType="application/json")
        print("combined_style.json 최초 저장 완료")
        return

    # 기존 스타일과 신규 스타일 LLM 병합
    prompt = (
        "아래 두 문서 스타일을 하나의 통합 스타일 가이드로 병합하세요.\n\n"
        "병합 원칙:\n"
        "- 공통된 특징은 하나로 통합\n"
        "- 서로 다른 특징은 '주로 X, 경우에 따라 Y' 형식으로 표현\n"
        "- 충돌하는 특징은 더 최근 문서(신규 스타일) 기준으로 결정\n"
        "- 원본 형식([문서 작성 스타일], [표현 규칙] 등 섹션 구조) 유지\n\n"
        "- 반드시 모든 섹션을 완성해서 출력할 것. 중간에 끊지 말 것\n\n"
        f"[기존 스타일]{newline}{existing}{newline}{newline}"
        f"[신규 스타일]{newline}n{new_style_desc}{newline}{newline}"
        "통합 스타일만 출력:"
    )
    resp = bedrock.invoke_model(
        modelId=MODEL_ID,
        body=json.dumps({
            "anthropic_version": "bedrock-2023-05-31",
            "max_tokens": 4000,
            "messages": [{"role": "user", "content": prompt}],
        }),
    )
    merged = json.loads(resp["body"].read())["content"][0]["text"].strip()

    body = json.dumps({"style_desc": merged}, ensure_ascii=False)
    s3_client.put_object(Bucket=S3_BUCKET, Key=key, Body=body,
                         ContentType="application/json")
    print("combined_style.json 병합 저장 완료")

        
        
def delete_template(user_id: str, template_id: str) -> bool:
    prefix = f"{S3_PREFIX}/{user_id}/{template_id}/"
    try:
        resp = s3_client.list_objects_v2(Bucket=S3_BUCKET, Prefix=prefix)
        objects = resp.get("Contents", [])
        if not objects:
            return True
        s3_client.delete_objects(
            Bucket=S3_BUCKET,
            Delete={"Objects": [{"Key": o["Key"]} for o in objects]},
        )
        print(f"템플릿 삭제 완료: {template_id} ({len(objects)}개 파일)")
        return True
    except Exception as e:
        print(f"템플릿 삭제 실패: {e}")
        return False
        
        
def _sanitize_messages(messages: list) -> list:
    def clean_text(text: str) -> str:
        text = "".join(c for c in str(text) if not (0xD800 <= ord(c) <= 0xDFFF))
        text = re.sub(r'[^\u0000-\uFFFF]', '', text)
        text = text.encode("utf-8", errors="ignore").decode("utf-8")
        return text

    return [
        {
            "role":    m["role"],
            "content": clean_text(m.get("content", "")),
        }
        for m in messages
    ]


def save_conversation(user_id: str, template_id: str, messages: list, session_id: str,
                      outline: str, user_mode: str = "", state: dict = None) -> None:
    key = f"{S3_PREFIX}/{user_id}/{template_id}/conversations/{session_id}.json"
    ts = datetime.now(tz=KST).strftime("%Y-%m-%d %H:%M")
    first_msg = next((m["content"] for m in messages if m["role"] == "user"), "새 대화")
    raw = {
        "session_id": session_id,
        "title":      first_msg[:10] + ("..." if len(first_msg) > 20 else ""),
        "saved_at":   ts,
        "messages":   _sanitize_messages(messages),
        "outline":    outline,
        "user_mode":  user_mode,
        "state":      state or {},          # ← 추가: 진행 상태 묶음
    }
    clean_data = json.loads(json.dumps(raw, ensure_ascii=True))
    body = json.dumps(clean_data, ensure_ascii=False).encode("utf-8")
    s3_client.put_object(Bucket=S3_BUCKET, Key=key, Body=body,
                         ContentType="application/json; charset=utf-8")
    print("대화 이력 저장 완료")


def load_conversation(user_id: str, template_id: str, session_id: str) -> dict:
    key = f"{S3_PREFIX}/{user_id}/{template_id}/conversations/{session_id}.json"
    try:
        obj  = s3_client.get_object(Bucket=S3_BUCKET, Key=key)
        return json.loads(obj["Body"].read())
    except Exception as e:
        if "NoSuchKey" in str(e) or "does not exist" in str(e):
            return {}
        print(f"  대화 로드 실패: {e}")
        return {}

def delete_conversation(user_id: str, template_id: str, session_id: str) -> bool:
    key = f"{S3_PREFIX}/{user_id}/{template_id}/conversations/{session_id}.json"
    try:
        s3_client.delete_object(Bucket=S3_BUCKET, Key=key)
        print(f"대화 삭제 완료: {key}")
        return True
    except Exception as e:
        print(f" 대화 삭제 실패: {e}")
        return False


def update_conversation_title(user_id: str, template_id: str, session_id: str, new_title: str):
    key = f"{S3_PREFIX}/{user_id}/{template_id}/conversations/{session_id}.json"
    try:
        obj  = s3_client.get_object(Bucket=S3_BUCKET, Key=key)
        data = json.loads(obj["Body"].read())
        data["title"] = new_title
        body = json.dumps(data, ensure_ascii=False).encode("utf-8")
        s3_client.put_object(Bucket=S3_BUCKET, Key=key, Body=body, ContentType="application/json; charset=utf-8")
    except Exception as e:
        print(f"제목 업데이트 실패: {e}")
        
        
def get_user_templates(user_id: str) -> list:
    try:
        prefix = f"{S3_PREFIX}/{user_id}/"
        resp = s3_client.list_objects_v2(
            Bucket=S3_BUCKET,
            Prefix=prefix,
            Delimiter="/",
        )
        templates = []
        for p in resp.get("CommonPrefixes", []):
            safe_tpl = p["Prefix"].rstrip("/").split("/")[-1]
            templates.append({"id": safe_tpl, "display": _read_display_name(user_id, safe_tpl)})
        return sorted(templates, key=lambda t: t["display"])
    except Exception as e:
        import traceback
        print(f"[get_user_templates] {traceback.format_exc()}")
        return []
        
# ──────────────────────────────────────────────────────────────
# AgentCore 저장 함수
#
# [네임스페이스 분리 원칙]
#   /writing/{actor_id}/    ← 문서 파일 분석 결과 (semantic)
#   /preference/{actor_id}/ ← 온보딩 응답 + !저장 키워드 피드백 (userPreference)
#
# [semantic 오염 차단]
#   save_style_to_agentcore(): 스타일 텍스트를 ASSISTANT role에 배치
#   → userPreferenceMemoryStrategy는 USER 발화만 스캔하므로
#     ASSISTANT role 텍스트에서 preference 추출 안 됨
# ──────────────────────────────────────────────────────────────


def save_style_to_agentcore(actor_id: str, style_desc: str):

    session_id = f"style-{actor_id}-{uuid.uuid4().hex}"
    ac_client.create_event(
        memoryId=MEMORY_ID,
        actorId=actor_id,
        sessionId=session_id,
        eventTimestamp=datetime.now(timezone(timedelta(hours=9))),
        payload=[
            {
                "conversational": {
                    "role": "USER",
                    "content": {"text": _sanitize(f"[문서 스타일 기록]\n{style_desc}")},
                }
            },
            {
                "conversational": {
                    "role": "ASSISTANT",
                    "content": {"text": "문서 스타일을 기록했습니다."},
                }
            },
        ],
    )
    print(f"   writing 저장 완료 → /writing/{actor_id}/")
    
    

def save_preference_to_agentcore(actor_id: str, pref_text: str):

    session_id = f"pref-{actor_id}-{uuid.uuid4().hex}"
    ac_client.create_event(
        memoryId=MEMORY_ID,
        actorId=actor_id,
        sessionId=session_id,
        eventTimestamp=datetime.now(timezone(timedelta(hours=9))),
        payload=[
            {
                "conversational": {
                    "role": "USER",
                    "content": {
                        "text": _sanitize(
                            f"나는 다음과 같은 문서 작성 스타일을 선호합니다:{newline}{pref_text}"
                        )
                    },
                }
            },
            {
                "conversational": {
                    "role": "ASSISTANT",
                    "content": {"text": "선호도를 기억하겠습니다."},
                }
            },
        ],
    )
    print(f"   preference 저장 완료 → /preference/{actor_id}/")


# ──────────────────────────────────────────────────────────────
# 이미지 텍스트 추출 (Bedrock Claude vision)
# ──────────────────────────────────────────────────────────────
def extract_text_from_image(image_path: str) -> str:

    ext = Path(image_path).suffix.lower()
    media_type_map = {
        ".jpg":  "image/jpeg",
        ".jpeg": "image/jpeg",
        ".png":  "image/png",
        ".gif":  "image/gif",
        ".webp": "image/webp",
    }
    media_type = media_type_map.get(ext)
    if not media_type:
        raise ValueError(f"지원하지 않는 이미지 형식: {ext}")

    with open(image_path, "rb") as f:
        image_data = base64.standard_b64encode(f.read()).decode("utf-8")

    resp = bedrock.invoke_model(
        modelId=MODEL_ID,
        body=json.dumps({
            "anthropic_version": "bedrock-2023-05-31",
            "max_tokens": 2000,
            "messages": [
                {
                    "role": "user",
                    "content": [
                        {
                            "type": "image",
                            "source": {
                                "type":       "base64",
                                "media_type": media_type,
                                "data":       image_data,
                            },
                        },
                        {
                            "type": "text",
                            "text": (
                                "이 이미지에서 모든 텍스트를 추출해주세요. "
                                "슬라이드 제목, 본문, bullet 항목, 표 내용 등 "
                                "보이는 텍스트를 원본 구조 그대로 출력하세요. "
                                "설명이나 부연은 하지 말고 텍스트만 출력하세요."
                            ),
                        },
                    ],
                }
            ],
        }),
    )
    extracted = json.loads(resp["body"].read())["content"][0]["text"]
    print(f"   이미지 텍스트 추출 완료 ({len(extracted)}자)")
    return extracted


def is_image_file(path: str) -> bool:
    return Path(path).suffix.lower() in IMAGE_EXTS


def is_doc_file(path: str) -> bool:
    return Path(path).suffix.lower() in DOC_EXTS



# ── 공통 헬퍼: x중심 비율로 좌/중/우 판정 ──
def _zone_by_ratio(cx: float, width: float) -> str:
    if not width or width <= 0:
        return "중"
    r = cx / width
    return "좌" if r < 0.4 else ("우" if r > 0.6 else "중")
    
# ──────────────────────────────────────────────────────────────
# PPTX 파싱
# ──────────────────────────────────────────────────────────────
def _extract_table_content(shape) -> dict:
    rows = [[cell.text.strip() for cell in row.cells] for row in shape.table.rows]
    rows = [r for r in rows if any(r)]
    if not rows:
        return {}
    return {
        "headers":   rows[0],
        "rows":      rows[1:],
        "row_count": len(rows) - 1,
        "col_count": len(rows[0]),
        "structure": (
            "비교형" if len(rows[0]) == 2 else
            "다중비교형" if len(rows[0]) >= 3 and len(rows) <= 5 else
            "목록형"
        ),
    }


def _find_slide_title(slide) -> str:
    """슬라이드 제목을 견고하게 찾는다.
    1순위: 제목 플레이스홀더(TITLE=1, CENTER_TITLE=0)
    2순위: 그 슬라이드에서 '가장 큰 폰트'의 텍스트 (절대 크기 아님 — 상대 최대).
           단 본문(그보다 작은 폰트)이 존재해 계층이 있을 때만. 크기 정보가
           전혀 없거나 모두 같은 크기면 제목을 특정하지 않는다(오인 방지)."""
    # 1) 제목 플레이스홀더
    for shape in slide.shapes:
        if shape.has_text_frame and shape.is_placeholder:
            try:
                if int(shape.placeholder_format.type) in (0, 1):
                    t = shape.text_frame.text.strip()
                    if t:
                        return t
            except Exception:
                pass
    # 2) 폴백: 문단별 (텍스트, 폰트pt) 수집
    items = []  # (text, size_pt)
    for shape in slide.shapes:
        if not shape.has_text_frame:
            continue
        for para in shape.text_frame.paragraphs:
            t = para.text.strip()
            if not t:
                continue
            sz = 0
            try:
                if para.runs and para.runs[0].font.size:
                    sz = para.runs[0].font.size.pt
                elif para.font.size:
                    sz = para.font.size.pt
            except Exception:
                sz = 0
            items.append((t, sz))
 
    sized = [(t, s) for t, s in items if s and s > 0]
    if not sized:
        return ""                      # 크기 정보 없음 → 제목 특정 불가
    max_sz = max(s for _, s in sized)
    body_szs = [s for _, s in sized if s < max_sz]
    if not body_szs:
        return ""                      # 전부 같은 크기 → 제목 계층 없음(오인 방지)
    # 가장 큰 폰트의 첫 텍스트를 제목으로 (문서마다 크기 달라도 상대적으로 최대면 제목)
    for t, s in sized:
        if s == max_sz:
            return t
    return ""
 
 
def _extract_pptx(pptx_path: str) -> dict:
    prs, slides_info = Presentation(pptx_path), []
    sw = prs.slide_width or 1
    for i, slide in enumerate(prs.slides, 1):
        sd = {"slide_num": i, "title": "", "bullets": [], "paragraphs": [],
              "tables": [], "content_items": []}
 
        # ── 제목 먼저 확정 (플레이스홀더 우선 → 폰트 폴백) ──
        sd["title"] = _find_slide_title(slide)
        _title_used = sd["title"]   # 본문에서 제외할 제목 텍스트
 
        for shape in slide.shapes:
            # 도형 x중심 → 좌/중/우
            try:
                cx = shape.left + (shape.width or 0) / 2 if shape.left is not None else None
                zone = _zone_by_ratio(cx, sw) if cx is not None else "중"
            except Exception:
                zone = "중"
 
            if shape.has_table:
                t = _extract_table_content(shape)
                if t:
                    sd["tables"].append(t)
                    sd["content_items"].append({"zone": zone, "text": "[표]"})
            elif shape.has_text_frame:
                for para in shape.text_frame.paragraphs:
                    text = para.text.strip()
                    if not text:
                        continue
                    if text == _title_used:      # 이미 제목으로 쓴 텍스트는 본문에서 제외
                        continue
                    if any(text.startswith(c) for c in ["•","-","▪","○","·"]) or para.level > 0:
                        sd["bullets"].append(text)
                    else:
                        sd["paragraphs"].append(text)
                    sd["content_items"].append({"zone": zone, "text": text})
        slides_info.append(sd)
        print(f"[제목확인] 슬라이드{i}: title='{sd['title']}' | 본문 {len(sd['content_items'])}개")
 
    bullet_cnt = sum(len(s["bullets"])    for s in slides_info)
    para_cnt   = sum(len(s["paragraphs"]) for s in slides_info)
    all_tables = [t for s in slides_info for t in s["tables"]]
    table_pats: dict = {}
    for t in all_tables:
        table_pats[t["structure"]] = table_pats.get(t["structure"], 0) + 1
 
    return {
        "slide_details": [
            {"slide_num": s["slide_num"], "title": s["title"],
             "content": s["bullets"] + s["paragraphs"], "tables": s["tables"],
             "content_items": s["content_items"]}
            for s in slides_info
            if s["title"] or s["bullets"] or s["paragraphs"] or s["tables"]
        ],
        "toc":            [s["title"] for s in slides_info if s["title"]],
        "slide_count":    len(prs.slides),
        "writing_style":  "개조식" if bullet_cnt > para_cnt else "서술형",
        "table_count":    len(all_tables),
        "table_patterns": table_pats,
    }
# ──────────────────────────────────────────────────────────────
# PDF 파싱 (pdfplumber)
# ──────────────────────────────────────────────────────────────

def _layout_from_spans(spans: list, page_w: float) -> str:
    if not spans:
        return "title_only"
    # 판정에 쓸 박스: 너무 넓은 것(전폭·배경)과 너무 좁은 것(제목·라벨·아이콘)은 제외
    core = [s for s in spans
            if page_w * 0.15 <= (s[1] - s[0]) < page_w * 0.70]
    n = len(core)
    if n == 0:
        return "single"
    band_lo, band_hi = page_w * 0.45, page_w * 0.55
    crossing   = [s for s in core if s[0] < band_hi and s[1] > band_lo]
    left_only  = any(r <= band_lo for l, r in core)
    right_only = any(l >= band_hi for l, r in core)
    if left_only and right_only and len(crossing) <= max(int(n * 0.25), 1):
        return "two_col"
    return "single"

def _extract_layout_from_page(page) -> dict:
    pw = page.width
    words = page.extract_words(extra_attrs=["size"]) or []
    font_sizes = [w.get("size", 0) for w in words if w.get("size")]
    title_size = max(font_sizes) if font_sizes else 0
 
    title = ""
    content_items = []
    for w in words:
        text = w["text"].strip()
        if not text:
            continue
        # 제목(가장 큰 글씨 첫 등장) 분리
        if w.get("size", 0) >= title_size * 0.95 and not title:
            title = text
            continue
        cx = (w["x0"] + w["x1"]) / 2
        content_items.append({"zone": _zone_by_ratio(cx, pw), "text": text})
 
    return {"title": title, "content_items": content_items}

def _extract_pdf(pdf_path: str) -> dict:
    import pdfplumber
    pages_info: list = []
 
    with pdfplumber.open(pdf_path) as pdf:
        for i, page in enumerate(pdf.pages, 1):
            pw = page.width or 1
            pd = {"slide_num": i, "title": "", "content_items": [], "tables": []}
 
            # ── 표 추출 (기존 유지) + 표 위치 태깅 ──
            table_bboxes = []
            try:
                found = page.find_tables()
                table_bboxes = [t.bbox for t in found]  # (x0,y0,x1,y1)
            except Exception:
                table_bboxes = []
 
            for ti, table in enumerate(page.extract_tables()):
                if not table:
                    continue
                rows = [[cell.strip() if cell else "" for cell in row] for row in table]
                rows = [r for r in rows if any(r)]
                if not rows:
                    continue
                pd["tables"].append({
                    "headers":   rows[0], "rows": rows[1:],
                    "row_count": len(rows) - 1, "col_count": len(rows[0]),
                    "structure": (
                        "비교형"     if len(rows[0]) == 2 else
                        "다중비교형"  if len(rows[0]) >= 3 and len(rows) <= 5 else
                        "목록형"
                    ),
                })
                # 표의 x중심으로 zone 태깅 (bbox가 있으면)
                if ti < len(table_bboxes):
                    bx0, _, bx1, _ = table_bboxes[ti]
                    cx = (bx0 + bx1) / 2
                    z = "좌" if cx / pw < 0.4 else ("우" if cx / pw > 0.6 else "중")
                else:
                    z = "중"
                pd["content_items"].append({"zone": z, "text": "[표]"})
 
            # ── 본문: 새 _extract_layout_from_page 에서 title + content_items(zone) ──
            layout = _extract_layout_from_page(page)
            pd["title"] = layout.get("title", "")
            pd["content_items"] += layout.get("content_items", [])
 
            pages_info.append(pd)
 
    # writing_style 계산용: bullet vs paragraph (content_items 텍스트에서 추정)
    def _is_bullet(t):
        return any(t.startswith(c) for c in ["•","-","▪","○","·","–"])
    bullet_cnt = sum(1 for p in pages_info for it in p["content_items"]
                     if it["text"] != "[표]" and _is_bullet(it["text"]))
    para_cnt   = sum(1 for p in pages_info for it in p["content_items"]
                     if it["text"] != "[표]" and not _is_bullet(it["text"]))
 
    all_tables = [t for p in pages_info for t in p["tables"]]
    table_pats: dict = {}
    for t in all_tables:
        table_pats[t["structure"]] = table_pats.get(t["structure"], 0) + 1
 
    return {
        "slide_details": [
            {"slide_num": p["slide_num"], "title": p["title"],
             "content": [it["text"] for it in p["content_items"] if it["text"] != "[표]"],
             "tables": p["tables"],
             "content_items": p["content_items"]}         # ← LLM 판정용 위치정보
            for p in pages_info
            if p["title"] or p["content_items"] or p["tables"]
        ],
        "toc":            [p["title"] for p in pages_info if p["title"]],
        "slide_count":    len(pages_info),
        "writing_style":  "개조식" if bullet_cnt > para_cnt else "서술형",
        "table_count":    len(all_tables),
        "table_patterns": table_pats,

    }


def extract_doc_style(file_path: str) -> dict:
    ext = Path(file_path).suffix.lower()
    if ext == ".pdf":
        return _extract_pdf(file_path)
    elif ext == ".pptx":
        return _extract_pptx(file_path)
    else:
        raise ValueError(f"지원하지 않는 파일 형식: {ext}")
           
        

def _read_display_name(user_id, safe_tpl):
    key = f"{S3_PREFIX}/{user_id}/{safe_tpl}/meta/template_meta.json"
    try:
        obj = s3_client.get_object(Bucket=S3_BUCKET, Key=key)
        return json.loads(obj["Body"].read())["display_name"]
    except Exception:
        return safe_tpl
        
        
        
def list_user_templates(user_id: str) -> list:
    if not user_id or not user_id.strip():
        return []
    prefix = f"{S3_PREFIX}/{user_id}/"
    try:
        resp = s3_client.list_objects_v2(
            Bucket=S3_BUCKET,
            Prefix=prefix,
            Delimiter="/",
        )
        templates = []
        for p in resp.get("CommonPrefixes", []):
            safe_tpl = p["Prefix"].rstrip("/").split("/")[-1]
            templates.append({"id": safe_tpl, "display": _read_display_name(user_id, safe_tpl)})
        return templates
    except Exception as e:
        print(f"템플릿 목록 조회 실패: {e}")
        return []
        
        
        
