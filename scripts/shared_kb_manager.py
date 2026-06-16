"""
shared_kb_manager.py

공용 Knowledge Base 파일 업로드 및 Ingestion 관리 스크립트.
챗봇 서버와 무관하게 관리자가 직접 CLI로 실행.

S3 경로 구조:
    shared/{folder}/raw/        ← 원본 업로드 파일
    shared/{folder}/processed/  ← Lambda 변환 결과 (Bedrock이 바라보는 경로)

사용 예시:
    # 파일 1개 업로드 + Ingestion
    python scripts/shared_kb_manager.py --action upload --folder policy --file ./docs/shared_kb_docs/policy/policy_2026.pdf

    # 파일 여러 개 업로드 + Ingestion
    python scripts/shared_kb_manager.py --action upload --folder policy \
        --file ./docs/shared_kb_docs/policy/policy_2026.pdf ./docs/shared_kb_docs/policy/policy_2025.pdf

    # 디렉토리 전체 업로드 + Ingestion
    python scripts/shared_kb_manager.py --action upload --folder policy --dir ./docs/shared_kb_docs/policy/

    # 특정 폴더 Ingestion만 실행 (S3에 이미 업로드된 파일 재처리)
    python scripts/shared_kb_manager.py --action ingest --folder policy

    # Ingestion 상태 확인
    python scripts/shared_kb_manager.py --action status --folder policy --job-id abc123

    # 등록된 폴더(Data Source) 목록 확인
    python scripts/shared_kb_manager.py --action list

    # 새 폴더(Data Source) 등록
    python scripts/shared_kb_manager.py --action add-folder --folder regulation

config/knowBase.yaml 의 shared_kb 섹션에 아래 항목이 채워져 있어야 합니다:
    shared_kb:
        kb_id:                   "your-shared-kb-id"
        kb_role_arn:             "arn:aws:iam::ACCOUNT:role/bedrock-kb-role"
        s3_bucket:               "your-s3-bucket"
        s3_intermediate_bucket:  "your-s3-intermediate-bucket"
        lambda_arn:              "arn:aws:lambda:ap-northeast-2:ACCOUNT:function:transform_lambda"
        embedding_model:         "amazon.titan-embed-text-v2:0"
        folders:
            policy:     "ds-id-policy"
            manual:     "ds-id-manual"
            notice:     "ds-id-notice"
        poll_interval:   5
        poll_timeout:    300
"""

from __future__ import annotations

import argparse
import hashlib
import io
import json
import re
import sys
import time
from pathlib import Path

import boto3
import pandas as pd
import yaml

# 프로젝트 루트를 sys.path에 추가 (scripts/ 에서 직접 실행하기 위함)
ROOT = Path(__file__).resolve().parent.parent
if str(ROOT) not in sys.path:
    sys.path.insert(0, str(ROOT))

# ──────────────────────────────────────────────
# 설정 로드 (.env 의 KB_* 변수 주입은 get_kb_config 가 처리)
# ──────────────────────────────────────────────
CONFIG_PATH = ROOT / "config" / "knowBase.yaml"

from wellbot.env import init_env  # noqa: E402
init_env()  # KB 모듈의 모듈레벨 os.getenv 보장 (다른 wellbot import 전에 호출)

from wellbot.services.knowledgebase.config import get_kb_config  # noqa: E402  (sys.path 설정 이후 import)

_config = get_kb_config()
_kb_cfg = _config["shared_kb"]

SHARED_KB_ID           = _kb_cfg["kb_id"]
KB_ROLE_ARN            = _kb_cfg["kb_role_arn"]
S3_BUCKET              = _kb_cfg["s3_bucket"]
S3_INTERMEDIATE_BUCKET = _kb_cfg["s3_intermediate_bucket"]
LAMBDA_ARN             = _kb_cfg["lambda_arn"]
EMBEDDING_MODEL        = _kb_cfg["embedding_model"]
POLL_INTERVAL          = _kb_cfg.get("poll_interval", 5)
POLL_TIMEOUT           = _kb_cfg.get("poll_timeout", 300)

# ──────────────────────────────────────────────
# 업로드 제한 설정
# ──────────────────────────────────────────────
ROWS_PER_SPLIT   = 50_000
TABULAR_EXTS     = {".xlsx", ".csv"}
CONVERTIBLE_EXTS = {".pptx"}     # Bedrock KB 미지원 → 업로드 전 json 변환
MAX_FILE_SIZES   = {
    ".txt":  30 * 1024 * 1024,
    ".md":   30 * 1024 * 1024,
    ".json": 30 * 1024 * 1024,
    ".csv":  None,
    ".xlsx": None,
}
MAX_FILE_SIZE_DEFAULT = 100 * 1024 * 1024  # 100MB

# ──────────────────────────────────────────────
# AWS 클라이언트
# ──────────────────────────────────────────────
_s3            = boto3.client("s3")
_bedrock_agent = boto3.client("bedrock-agent", region_name="ap-northeast-2")


# ──────────────────────────────────────────────
# S3 경로 헬퍼
# ──────────────────────────────────────────────
def _raw_prefix(folder: str) -> str:
    return f"shared/{folder}/raw/"

def _processed_prefix(folder: str) -> str:
    return f"shared/{folder}/processed/"


def _ds_name(folder: str) -> str:
    """Bedrock 데이터소스 이름 생성 (ASCII 안전).

    Bedrock 리소스 이름은 영문/숫자/하이픈/언더스코어만 허용하므로 한글 폴더명을
    그대로 쓸 수 없다. ASCII 슬러그 + 폴더명 해시 8자로 고유하고 유효한 이름을 만든다.
    한글 등 비ASCII 폴더는 슬러그가 비므로 해시만 사용 (콘솔 식별은 description 의
    한글 폴더명으로 한다). 영문 폴더는 슬러그가 남아 콘솔에서도 읽기 쉽다.
        'policy' → 'aiinno-bedrock-kb-ds-shared-policy-1a2b3c4d'
        '규정'   → 'aiinno-bedrock-kb-ds-shared-7e8f9a0b'
    """
    ascii_slug = re.sub(r"[^a-zA-Z0-9]+", "-", folder).strip("-").lower()
    digest = hashlib.md5(folder.encode("utf-8")).hexdigest()[:8]
    suffix = f"{ascii_slug}-{digest}" if ascii_slug else digest
    return f"aiinno-bedrock-kb-ds-shared-{suffix}"[:100]


# ──────────────────────────────────────────────
# config/knowBase.yaml 폴더 → Data Source 매핑 관리
# ──────────────────────────────────────────────
def _get_folders() -> dict:
    return _kb_cfg.get("folders", {})


def _get_data_source_id(folder: str) -> str:
    folders = _get_folders()
    if folder not in folders:
        raise ValueError(
            f"폴더 '{folder}'가 등록되어 있지 않습니다. "
            f"등록된 폴더: {list(folders.keys())}\n"
            f"새 폴더 등록: python scripts/shared_kb_manager.py --action add-folder --folder {folder}"
        )
    return folders[folder]


def _save_data_source_id(folder: str, data_source_id: str) -> None:
    """
    config/knowBase.yaml의 shared_kb.folders에 폴더를 추가.
    yaml.dump 대신 텍스트 삽입으로 기존 주석/형식을 보존.
    """
    content = CONFIG_PATH.read_text(encoding="utf-8")

    new_entry = f"    {folder}: \"{data_source_id}\""

    # folders: {} (빈 dict) → folders:\n    folder: "ds-id" 로 변환
    if "folders: {}" in content:
        content = content.replace(
            "folders: {}",
            f"folders:\n{new_entry}",
        )
    elif "folders:" in content:
        lines = content.split("\n")
        insert_idx = None
        in_folders = False
        for i, line in enumerate(lines):
            stripped = line.strip()
            if stripped.startswith("folders:"):
                in_folders = True
                insert_idx = i + 1
                continue
            if in_folders:
                if line.startswith("    ") and (stripped.startswith("#") or ":" in stripped):
                    insert_idx = i + 1
                else:
                    break

        if insert_idx is not None:
            lines.insert(insert_idx, new_entry)
            content = "\n".join(lines)
    else:
        print("[Config] folders 키를 찾을 수 없어 yaml.dump로 fallback")
        with open(CONFIG_PATH, "r", encoding="utf-8") as f:
            full_config = yaml.safe_load(f)
        full_config["shared_kb"].setdefault("folders", {})[folder] = data_source_id
        with open(CONFIG_PATH, "w", encoding="utf-8") as f:
            yaml.dump(full_config, f, allow_unicode=True, default_flow_style=False)
        _kb_cfg["folders"] = full_config["shared_kb"]["folders"]
        return

    CONFIG_PATH.write_text(content, encoding="utf-8")

    # 메모리 캐시도 업데이트
    _kb_cfg.setdefault("folders", {})[folder] = data_source_id
    print(f"[Config] 폴더 등록 완료: {folder} → {data_source_id}")


# ──────────────────────────────────────────────
# Data Source 생성 (새 폴더 등록 시)
# ──────────────────────────────────────────────
def add_folder(folder: str) -> str:
    folders = _get_folders()
    if folder in folders:
        print(f"[Folder] 이미 등록된 폴더: {folder} → {folders[folder]}")
        return folders[folder]

    print(f"[Folder] 새 폴더 Data Source 생성: folder={folder}")

    resp = _bedrock_agent.create_data_source(
        knowledgeBaseId=SHARED_KB_ID,
        # 이름은 ASCII 슬러그+해시 (한글 폴더 지원). 한글 폴더명은 description 으로 식별.
        name=_ds_name(folder),
        description=f"공용 KB - {folder} 폴더",
        dataSourceConfiguration={
            "type": "S3",
            "s3Configuration": {
                "bucketArn":         f"arn:aws:s3:::{S3_BUCKET}",
                "inclusionPrefixes": [_raw_prefix(folder)],
            },
        },
        vectorIngestionConfiguration={
            "chunkingConfiguration": {
                "chunkingStrategy": "NONE",
            },
            "customTransformationConfiguration": {
                "intermediateStorage": {
                    "s3Location": {
                        "uri": f"s3://{S3_INTERMEDIATE_BUCKET}/{_processed_prefix(folder)}",
                    },
                },
                "transformations": [{
                    "stepToApply": "POST_CHUNKING",
                    "transformationFunction": {
                        "transformationLambdaConfiguration": {
                            "lambdaArn": LAMBDA_ARN,
                        },
                    },
                }],
            },
        },
    )
    data_source_id = resp["dataSource"]["dataSourceId"]
    _save_data_source_id(folder, data_source_id)
    return data_source_id


# ──────────────────────────────────────────────
# 파일 크기 검증
# ──────────────────────────────────────────────
def _validate_file_size(file_path: str) -> None:
    path  = Path(file_path)
    ext   = path.suffix.lower()
    size  = path.stat().st_size
    limit = MAX_FILE_SIZES.get(ext, MAX_FILE_SIZE_DEFAULT)
    if limit is None:
        return
    if size > limit:
        limit_mb  = limit // (1024 * 1024)
        actual_mb = size / (1024 * 1024)
        raise ValueError(
            f"파일 크기 초과: {path.name} ({actual_mb:.1f}MB). "
            f"{ext} 파일은 {limit_mb}MB 이하만 업로드 가능합니다."
        )


# ──────────────────────────────────────────────
# 분할 업로드 (xlsx / csv)
# ──────────────────────────────────────────────
def _cleanup_existing_parts(folder: str, stem: str, ext: str) -> None:
    """재업로드 시 파트 수가 달라져서 오래된 파트가 남는 것을 방지."""
    prefix    = _raw_prefix(folder)
    paginator = _s3.get_paginator("list_objects_v2")
    deleted   = 0
    for page in paginator.paginate(Bucket=S3_BUCKET, Prefix=prefix):
        for obj in page.get("Contents", []):
            key      = obj["Key"]
            filename = key.split("/")[-1]
            if filename.startswith(f"{stem}_part") and filename.endswith(ext):
                _s3.delete_object(Bucket=S3_BUCKET, Key=key)
                deleted += 1
                print(f"[S3] 기존 파트 삭제: {key}")
    if deleted:
        print(f"[S3] 기존 파트 {deleted}개 정리 완료: stem={stem}, ext={ext}")


def _split_and_upload_tabular(folder: str, file_path: str) -> list[str]:
    """xlsx/csv를 ROWS_PER_SPLIT 행 단위로 분할해서 S3 raw/ 에 저장."""
    path = Path(file_path)
    ext  = path.suffix.lower()
    stem = path.stem

    _cleanup_existing_parts(folder, stem, ext)

    if ext == ".csv":
        try:
            df = pd.read_csv(file_path)
        except UnicodeDecodeError:
            df = pd.read_csv(file_path, encoding="cp949")
    else:
        df = pd.read_excel(file_path)

    total_rows = len(df)
    print(f"[Upload] 분할 업로드 시작: {path.name}, 총 {total_rows}행, {ROWS_PER_SPLIT}행씩 분할")

    uris = []
    for i, start in enumerate(range(0, total_rows, ROWS_PER_SPLIT)):
        chunk_df       = df.iloc[start:start + ROWS_PER_SPLIT]
        split_filename = f"{stem}_part{i + 1}{ext}"
        buf            = io.BytesIO()

        if ext == ".csv":
            chunk_df.to_csv(buf, index=False)
        else:
            chunk_df.to_excel(buf, index=False)

        buf.seek(0)
        key = f"{_raw_prefix(folder)}{split_filename}"
        _s3.put_object(Bucket=S3_BUCKET, Key=key, Body=buf.read())
        uri = f"s3://{S3_BUCKET}/{key}"
        print(
            f"[S3] 분할 업로드: {uri} "
            f"(rows {start}~{min(start + ROWS_PER_SPLIT, total_rows) - 1})"
        )
        uris.append(uri)

    print(f"[Upload] 분할 업로드 완료: {path.name} → {len(uris)}개 파트")
    return uris


# ──────────────────────────────────────────────
# pptx → json 변환
# ──────────────────────────────────────────────
def convert_pptx_to_json(file_path: str) -> str:
    """
    pptx 파일을 슬라이드별 구조화된 JSON으로 변환.
    Bedrock KB가 pptx를 직접 지원하지 않으므로 업로드 전에 변환.
    반환: 변환된 json 파일 경로 (예: report.pptx → report_pptx.json)
    """
    from pptx import Presentation

    path = Path(file_path)
    prs  = Presentation(str(path))
    result = {}

    for slide_num, slide in enumerate(prs.slides, 1):
        slide_title = "No Title"
        if slide.shapes.title and slide.shapes.title.has_text_frame:
            slide_title = slide.shapes.title.text.strip() or "No Title"

        parts = []

        for shape in slide.shapes:
            if shape.has_table:
                table = shape.table
                rows  = [[cell.text.strip() for cell in row.cells] for row in table.rows]
                if rows:
                    headers = rows[0]
                    for row in rows[1:]:
                        row_text = " | ".join(
                            f"{headers[i]}: {cell}" for i, cell in enumerate(row)
                            if i < len(headers)
                        )
                        if row_text.strip():
                            parts.append(row_text)
            elif shape.has_text_frame:
                for para in shape.text_frame.paragraphs:
                    text = para.text.strip()
                    if text:
                        parts.append(text)

        if slide.has_notes_slide:
            notes = slide.notes_slide.notes_text_frame.text.strip()
            if notes:
                parts.append(f"[노트] {notes}")

        if parts:
            key = f"slide_{slide_num}_{slide_title}"
            result[key] = "\n".join(parts)

    json_path = path.parent / f"{path.stem}_pptx.json"
    json_path.write_text(
        json.dumps(result, ensure_ascii=False, indent=2),
        encoding="utf-8",
    )
    print(f"[Convert] pptx → json 변환: {path.name} → {json_path.name} ({len(result)}슬라이드)")
    return str(json_path)


# ──────────────────────────────────────────────
# 파일 수집 유틸
# ──────────────────────────────────────────────
def collect_files_from_dir(dir_path: str) -> list[str]:
    supported = {".pdf", ".docx", ".pptx", ".xlsx", ".csv", ".md", ".txt", ".json", ".html", ".htm"}
    paths = [
        str(p) for p in Path(dir_path).iterdir()
        if p.is_file() and p.suffix.lower() in supported
    ]
    print(f"[Dir] {len(paths)}개 파일 수집: {dir_path}")
    return paths


# ──────────────────────────────────────────────
# 파일 업로드
# ──────────────────────────────────────────────
def upload_files(folder: str, file_paths: list[str]) -> list[str]:
    """
    여러 파일을 S3 shared/{folder}/raw/ 에 업로드.
    - pptx: 원본을 S3에 보관 후 json으로 변환하여 인덱싱용 업로드
    - xlsx/csv: ROWS_PER_SPLIT 행 단위로 분할 업로드
    - 그 외: 형식별 크기 제한 검증 후 단일 업로드 (기본 100MB)
    - 업로드 실패 시 부분 업로드된 파일 롤백
    반환: 업로드된 S3 URI 목록
    """
    folders = _get_folders()
    if folder not in folders:
        print(f"[Upload] 미등록 폴더 → 자동 생성: {folder}")
        add_folder(folder)

    for file_path in file_paths:
        path = Path(file_path)
        if not path.exists():
            raise FileNotFoundError(f"파일 없음: {file_path}")
        _validate_file_size(file_path)

    # pptx → json 변환 (원본은 S3에 별도 보관, 인덱싱은 json으로)
    converted_files: list[str] = []
    resolved_paths: list[str] = []
    for file_path in file_paths:
        ext = Path(file_path).suffix.lower()
        if ext in CONVERTIBLE_EXTS:
            # 원본 pptx를 S3에 보관 (다운로드용)
            original_key = f"{_raw_prefix(folder)}{Path(file_path).name}"
            with open(file_path, "rb") as f:
                _s3.put_object(Bucket=S3_BUCKET, Key=original_key, Body=f)
            print(f"[S3] 원본 보관: s3://{S3_BUCKET}/{original_key}")
            # json 변환본 (Bedrock 인덱싱용)
            json_path = convert_pptx_to_json(file_path)
            resolved_paths.append(json_path)
            converted_files.append(json_path)
        else:
            resolved_paths.append(file_path)

    uploaded_uris: list[str] = []
    try:
        for file_path in resolved_paths:
            path = Path(file_path)
            ext  = path.suffix.lower()

            if ext in TABULAR_EXTS:
                uris = _split_and_upload_tabular(folder, file_path)
            else:
                key = f"{_raw_prefix(folder)}{path.name}"
                with open(path, "rb") as f:
                    _s3.put_object(Bucket=S3_BUCKET, Key=key, Body=f)
                uri  = f"s3://{S3_BUCKET}/{key}"
                print(f"[S3] 업로드 완료: {uri}")
                uris = [uri]

            uploaded_uris.extend(uris)

    except Exception:
        if uploaded_uris:
            print(f"[S3] 업로드 실패, 롤백 시작: {len(uploaded_uris)}개 삭제")
            for uri in uploaded_uris:
                key = uri.replace(f"s3://{S3_BUCKET}/", "")
                try:
                    _s3.delete_object(Bucket=S3_BUCKET, Key=key)
                except Exception as del_err:
                    print(f"[S3] 롤백 실패: {key}, {del_err}")
        raise
    finally:
        for tmp in converted_files:
            try:
                Path(tmp).unlink(missing_ok=True)
            except Exception:
                pass

    return uploaded_uris


# ──────────────────────────────────────────────
# Ingestion
# ──────────────────────────────────────────────
def start_ingestion(folder: str) -> str:
    data_source_id = _get_data_source_id(folder)
    resp = _bedrock_agent.start_ingestion_job(
        knowledgeBaseId=SHARED_KB_ID,
        dataSourceId=data_source_id,
    )
    job_id = resp["ingestionJob"]["ingestionJobId"]
    print(f"[Bedrock] Ingestion 시작: folder={folder}, job_id={job_id}")
    return job_id


def poll_ingestion_status(folder: str, job_id: str) -> str:
    data_source_id = _get_data_source_id(folder)
    start = time.time()

    while time.time() - start < POLL_TIMEOUT:
        resp   = _bedrock_agent.get_ingestion_job(
            knowledgeBaseId=SHARED_KB_ID,
            dataSourceId=data_source_id,
            ingestionJobId=job_id,
        )
        status = resp["ingestionJob"]["status"]
        stats  = resp["ingestionJob"].get("statistics", {})
        print(
            f"[Bedrock] Ingestion 상태: folder={folder}, status={status}, "
            f"scanned={stats.get('numberOfDocumentsScanned', '-')}, "
            f"indexed={stats.get('numberOfNewDocumentsIndexed', '-')}, "
            f"failed={stats.get('numberOfDocumentsFailed', '-')}"
        )
        if status in ("COMPLETE", "FAILED", "STOPPED"):
            return status
        time.sleep(POLL_INTERVAL)

    raise TimeoutError(f"Ingestion 완료 대기 타임아웃: job_id={job_id}")


def upload_and_ingest(folder: str, file_paths: list[str]) -> None:
    """파일 목록 전체 업로드 → Ingestion 1회 실행 → 완료 대기."""
    uploaded = upload_files(folder, file_paths)
    if not uploaded:
        print("업로드된 파일이 없습니다.")
        return

    total_files = len(file_paths)
    total_uris  = len(uploaded)
    print(f"\n📤 {total_files}개 파일 업로드 완료 (S3 오브젝트 {total_uris}개)")
    job_id = start_ingestion(folder)
    print(f"⚙️  Ingestion 시작: job_id={job_id}")
    status = poll_ingestion_status(folder, job_id)
    print(f"✅ 완료: folder={folder}, status={status}")


# ──────────────────────────────────────────────
# 목록 출력
# ──────────────────────────────────────────────
def list_folders() -> None:
    folders = _get_folders()
    if not folders:
        print("등록된 폴더가 없습니다.")
        return
    print(f"\n{'폴더':<20} {'Data Source ID'}")
    print("-" * 60)
    for folder, ds_id in folders.items():
        print(f"{folder:<20} {ds_id}")


# ──────────────────────────────────────────────
# CLI 진입점
# ──────────────────────────────────────────────
def _parse_args():
    parser = argparse.ArgumentParser(description="공용 KB 관리 스크립트")
    parser.add_argument(
        "--action",
        required=True,
        choices=["upload", "ingest", "status", "list", "add-folder"],
    )
    parser.add_argument("--folder",  help="폴더명 (예: policy, manual, notice)")
    parser.add_argument("--file",    nargs="+", help="업로드할 파일 경로")
    parser.add_argument("--dir",     help="업로드할 디렉토리 경로")
    parser.add_argument("--job-id",  help="Ingestion Job ID (status 확인용)")
    return parser.parse_args()


if __name__ == "__main__":
    args = _parse_args()

    if args.action == "list":
        list_folders()

    elif args.action == "add-folder":
        if not args.folder:
            print("--folder 옵션이 필요합니다.")
            sys.exit(1)
        ds_id = add_folder(args.folder)
        print(f"✅ 폴더 등록 완료: {args.folder} → {ds_id}")

    elif args.action == "upload":
        if not args.folder:
            print("--folder 옵션이 필요합니다.")
            sys.exit(1)
        if not args.file and not args.dir:
            print("--file 또는 --dir 옵션이 필요합니다.")
            sys.exit(1)

        file_paths = []
        if args.dir:
            file_paths.extend(collect_files_from_dir(args.dir))
        if args.file:
            file_paths.extend(args.file)

        if not file_paths:
            print("업로드할 파일이 없습니다.")
            sys.exit(1)

        print(f"📂 업로드 대상: {len(file_paths)}개 파일 → folder={args.folder}")
        upload_and_ingest(args.folder, file_paths)

    elif args.action == "ingest":
        if not args.folder:
            print("--folder 옵션이 필요합니다.")
            sys.exit(1)
        job_id = start_ingestion(args.folder)
        print(f"Ingestion 시작: job_id={job_id}")
        status = poll_ingestion_status(args.folder, job_id)
        print(f"✅ 완료: status={status}")

    elif args.action == "status":
        if not args.folder or not args.job_id:
            print("--folder, --job-id 옵션이 필요합니다.")
            sys.exit(1)
        status = poll_ingestion_status(args.folder, args.job_id)
        print(f"상태: {status}")
