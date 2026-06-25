"""
Bedrock KB Custom Transformation Lambda

Bedrock Ingestion Job 이 실행될 때 자동 호출.
- 입력: S3 raw/ 의 파일 경로 (Bedrock 이 event 로 전달)
- 처리: 파일 형식별 파싱 + 청킹
- 출력: S3 processed/ 에 Bedrock KB 포맷 JSON 저장 후 경로 반환

Bedrock KB contentBody 포맷:
{
    "fileContents": [
        {
            "contentBody": "청킹된 텍스트",
            "contentType": "TEXT",
            "contentMetadata": { "key": "value" }
        },
        ...
    ]
}

Requirements (Lambda Layer):
    pdfplumber, python-docx, python-pptx, openpyxl, pandas
"""

import io
import json
import logging
import os
import re
import uuid
from dataclasses import dataclass, field
from typing import Any, Dict, List, Optional

import boto3
import docx
import pandas as pd
import pdfplumber
from bs4 import BeautifulSoup
from pptx import Presentation

logger = logging.getLogger()
logger.setLevel(logging.INFO)

s3 = boto3.client("s3")

# ──────────────────────────────────────────────
# 환경변수 기반 설정
# ──────────────────────────────────────────────
CHUNKER_TYPE    = os.environ.get("CHUNKER_TYPE", "recursive")       # "fixed" | "recursive"
INTERMEDIATE_BUCKET = os.environ["INTERMEDIATE_BUCKET"]
CHUNK_SIZE      = int(os.environ.get("CHUNK_SIZE", "1500"))
CHUNK_OVERLAP   = int(os.environ.get("CHUNK_OVERLAP", "200"))
ROWS_PER_CHUNK  = int(os.environ.get("ROWS_PER_CHUNK", "15"))             # 기본값 (csv/xlsx/docx 표/pptx 표)
PDF_TABLE_ROWS_PER_CHUNK = int(os.environ.get("PDF_TABLE_ROWS_PER_CHUNK", "5"))  # PDF 표 전용 (정밀 검색용)
MD_CHUNK_SIZE   = int(os.environ.get("MD_CHUNK_SIZE", "2000"))
JSON_CHUNK_SIZE = int(os.environ.get("JSON_CHUNK_SIZE", "1200"))


logger.info(
    f"[Config] chunker={CHUNKER_TYPE}, chunk_size={CHUNK_SIZE}, "
    f"overlap={CHUNK_OVERLAP}, rows_per_chunk={ROWS_PER_CHUNK}, "
    f"pdf_table_rows_per_chunk={PDF_TABLE_ROWS_PER_CHUNK}"
)


# ──────────────────────────────────────────────
# 데이터 구조
# ──────────────────────────────────────────────
@dataclass
class Chunk:
    text: str
    metadata: Dict[str, Any] = field(default_factory=dict)
    chunk_id: str = field(default_factory=lambda: str(uuid.uuid4()))


# ──────────────────────────────────────────────
# S3 유틸
# ──────────────────────────────────────────────
def _parse_s3_uri(uri: str):
    """s3://bucket/key → (bucket, key)"""
    uri = uri.replace("s3://", "")
    bucket, key = uri.split("/", 1)
    return bucket, key


def _read_s3_bytes(bucket: str, key: str) -> bytes:
    resp = s3.get_object(Bucket=bucket, Key=key)
    return resp["Body"].read()


MAX_METADATA_BYTES = 2048  # S3 Vectors filterable metadata 크기 제한


def _trim_metadata(meta: Dict[str, Any], limit: int = MAX_METADATA_BYTES) -> Dict[str, Any]:
    """
    메타데이터를 JSON 직렬화 기준 limit 바이트 이하로 축소
    1차: columns 필드를 컬럼 수 요약으로 대체
    2차: source 경로를 파일명만으로 축약
    3차: 그래도 초과하면 가장 큰 값부터 제거
    """
    def _byte_size(d: Dict) -> int:
        return len(json.dumps(d, ensure_ascii=False).encode("utf-8"))

    trimmed = dict(meta)

    # 1차: columns 리스트 → 컬럼 수 요약
    if "columns" in trimmed and isinstance(trimmed["columns"], list):
        col_count = len(trimmed["columns"])
        trimmed["columns"] = f"{col_count}개 컬럼"

    if _byte_size(trimmed) <= limit:
        return trimmed

    # 2차: source 경로 → 파일명만
    if "source" in trimmed and "/" in str(trimmed["source"]):
        trimmed["source"] = str(trimmed["source"]).rsplit("/", 1)[-1]

    if _byte_size(trimmed) <= limit:
        return trimmed

    # 3차: 큰 값부터 제거 (필수 키 보존)
    essential = {"doc_type", "chunk_index", "page", "sheet", "slide"}
    removable = sorted(
        [(k, len(json.dumps(v, ensure_ascii=False).encode("utf-8")))
         for k, v in trimmed.items() if k not in essential],
        key=lambda x: -x[1],
    )
    for k, _ in removable:
        del trimmed[k]
        if _byte_size(trimmed) <= limit:
            break

    if _byte_size(trimmed) > limit:
        logger.warning(f"[Meta] 축소 후에도 {_byte_size(trimmed)}B > {limit}B, 필수 키만 유지")

    return trimmed


def _write_chunks_to_s3(bucket: str, key: str, chunks: List[Chunk]) -> str:
    """
    청크 리스트를 Bedrock KB contentBody 포맷 JSON 으로 S3 에 저장
    저장 경로: processed/ prefix 로 변환하여 intermediate bucket 에 저장
    반환: 저장된 S3 URI
    """
    file_contents = [
        {
            "contentBody": c.text,
            "contentType": "TEXT",
            "contentMetadata": _trim_metadata(c.metadata),
        }
        for c in chunks
        if c.text.strip()
    ]
    payload = json.dumps({"fileContents": file_contents}, ensure_ascii=False)

    # raw/ → processed/ 경로 변환, 확장자를 .json 으로 변경
    processed_key = re.sub(r"^(.*/)raw/", r"\1processed/", key)
    processed_key = re.sub(r"\.[^.]+$", ".json", processed_key)

    s3.put_object(
        Bucket=bucket,
        Key=processed_key,
        Body=payload.encode("utf-8"),
        ContentType="application/json",
    )
    logger.info(f"[S3] 저장 완료: s3://{bucket}/{processed_key} ({len(file_contents)}청크)")
    return f"s3://{bucket}/{processed_key}"


# ──────────────────────────────────────────────
# 공통 텍스트 청킹 유틸
# ──────────────────────────────────────────────
_RECURSIVE_SEPARATORS = ["\n\n", "\n", ". ", " ", ""]


def _recursive_split(text: str, separators: List[str], chunk_size: int) -> List[str]:
    """구분자 우선순위에 따라 재귀적으로 분할"""
    if not separators:
        return [text[i:i + chunk_size] for i in range(0, len(text), chunk_size)]

    sep = separators[0]
    parts = text.split(sep) if sep else list(text)

    results = []
    for part in parts:
        if not part:
            continue
        if len(part) <= chunk_size:
            results.append(part)
        else:
            results.extend(_recursive_split(part, separators[1:], chunk_size))
    return results


def _recursive_merge(parts: List[str], chunk_size: int, overlap: int) -> List[str]:
    """분할된 조각을 chunk_size 이하로 합치되 overlap 적용"""
    merged, current = [], ""
    for part in parts:
        candidate = (current + " " + part).strip() if current else part
        if len(candidate) <= chunk_size:
            current = candidate
        else:
            if current:
                merged.append(current)
            overlap_text = current[-overlap:] if overlap else ""
            current = (overlap_text + " " + part).strip() if overlap_text else part
    if current:
        merged.append(current)
    return merged


def _chunk_text(
    text: str,
    source: str,
    doc_type: str,
    chunk_size: Optional[int] = None,
    overlap: Optional[int] = None,
    extra_meta: Optional[Dict] = None,
) -> List[Chunk]:
    """
    환경변수 CHUNKER_TYPE 에 따라 fixed 또는 recursive 청킹 수행
    chunk_size/overlap 미지정 시 환경변수 기본값 사용
    """
    cs = max(chunk_size or CHUNK_SIZE, 1)
    ov = overlap or CHUNK_OVERLAP
    # overlap 이 chunk_size 이상이면 fixed 모드의 step(cs-ov)이 0/음수가 되어
    # while 루프가 진행되지 않음(무한 루프 → Lambda 타임아웃). 0~cs-1 로 클램프.
    ov = max(0, min(ov, cs - 1))

    if CHUNKER_TYPE == "recursive":
        raw_parts = _recursive_split(text, _RECURSIVE_SEPARATORS, cs)
        merged = _recursive_merge(raw_parts, cs, ov)
        chunks = []
        for idx, part in enumerate(merged):
            part = part.strip()
            if not part:
                continue
            meta = {
                "source": source,
                "doc_type": doc_type,
                "chunk_index": idx,
                "chunker": "recursive",
            }
            if extra_meta:
                meta.update(extra_meta)
            chunks.append(Chunk(text=part, metadata=meta))
        return chunks
    else:
        # fixed-size 슬라이딩 윈도우
        chunks, start, idx = [], 0, 0
        while start < len(text):
            end = start + cs
            part = text[start:end].strip()
            if part:
                meta = {
                    "source": source,
                    "doc_type": doc_type,
                    "chunk_index": idx,
                    "start_char": start,
                    "chunker": "fixed",
                }
                if extra_meta:
                    meta.update(extra_meta)
                chunks.append(Chunk(text=part, metadata=meta))
                idx += 1
            start += cs - ov
        return chunks


def _chunk_table_rows(
    headers: List[str],
    rows: List[List[str]],
    source: str,
    doc_type: str,
    rows_per_chunk: Optional[int] = None,
    extra_meta: Optional[Dict] = None,
) -> List[Chunk]:
    """
    헤더 + 행 조합으로 텍스트 변환 후 청킹
    rows_per_chunk 미지정 시 환경변수 기본값 사용

    docx/pptx 의 병합 셀, 불규칙 테이블 등으로 헤더 < 데이터 행 길이가 되는
    경우를 방어: 부족한 헤더는 'col{N}' 으로 자동 보강
    """
    rpc = rows_per_chunk or ROWS_PER_CHUNK

    # 헤더보다 더 긴 row 가 있으면 헤더를 확장하여 IndexError 방지
    max_cols = max([len(headers)] + [len(r) for r in rows]) if rows else len(headers)
    headers_ext = list(headers) + [
        f"col{i}" for i in range(len(headers), max_cols)
    ]

    chunks = []
    for start in range(0, len(rows), rpc):
        batch = rows[start:start + rpc]
        rows_text = "\n".join(
            " | ".join(
                f"{headers_ext[i]}: {cell}" for i, cell in enumerate(row)
            )
            for row in batch
        )
        if rows_text.strip():
            meta = {
                "source":    source,
                "doc_type":  doc_type,
                "row_start": start,
                "row_end":   min(start + rpc, len(rows)) - 1,
                "columns":   headers_ext,
            }
            if extra_meta:
                meta.update(extra_meta)
            chunks.append(Chunk(text=rows_text, metadata=meta))
    return chunks


# ──────────────────────────────────────────────
# 파일 형식별 파서
# ──────────────────────────────────────────────
def parse_html(data: bytes, source: str) -> List[Chunk]:
    """
    BeautifulSoup 으로 HTML 태그 제거 후 텍스트 추출
    - <script>, <style> 태그는 제거
    - 제목(<title>)을 첫 줄에 포함
    """
    soup = BeautifulSoup(data.decode("utf-8", errors="replace"), "html.parser")

    # script, style 제거
    for tag in soup(["script", "style"]):
        tag.decompose()

    title = soup.title.string.strip() if soup.title and soup.title.string else ""
    text  = soup.get_text(separator="\n", strip=True)

    if title:
        text = f"{title}\n\n{text}"

    return _chunk_text(text, source, "html")


def parse_txt(data: bytes, source: str) -> List[Chunk]:
    text = data.decode("utf-8", errors="replace")
    return _chunk_text(text, source, "txt")


def parse_md(data: bytes, source: str) -> List[Chunk]:
    """헤더 기준 섹션 분리 후 청킹 (chromadb_chunking.MarkdownChunker 동일 로직)"""
    content  = data.decode("utf-8", errors="replace")
    sections = re.split(r"(?=\n#{1,3} )", content)
    chunks   = []
    for idx, section in enumerate(sections):
        section = section.strip()
        if not section:
            continue
        header_match = re.match(r"^(#{1,3})\s+(.+)", section)
        header = header_match.group(2) if header_match else "No Header"
        level  = len(header_match.group(1)) if header_match else 0
        for c in _chunk_text(section, source, "markdown", chunk_size=MD_CHUNK_SIZE, extra_meta={
            "section_header": header,
            "header_level":   level,
            "section_index":  idx,
        }):
            chunks.append(c)
    return chunks


def parse_json(data: bytes, source: str) -> List[Chunk]:
    """최상위 키 단위 분할 (chromadb_chunking.JsonChunker 동일 로직)"""
    content = json.loads(data.decode("utf-8"))
    chunks  = []
    for idx, (key, value) in enumerate(content.items()):
        text = f"{key}: {json.dumps(value, ensure_ascii=False, indent=2)}"
        for c in _chunk_text(text, source, "json", chunk_size=JSON_CHUNK_SIZE, extra_meta={
            "top_key":   key,
            "key_index": idx,
        }):
            chunks.append(c)
    return chunks


def parse_csv(data: bytes, source: str) -> List[Chunk]:
    df = pd.read_csv(io.BytesIO(data))
    headers = list(df.columns)
    rows    = [[str(v) for v in row] for _, row in df.iterrows()]
    return _chunk_table_rows(headers, rows, source, "csv")


def parse_xlsx(data: bytes, source: str) -> List[Chunk]:
    """시트별로 처리, 시트명을 메타데이터에 포함"""
    xl     = pd.ExcelFile(io.BytesIO(data))
    chunks = []
    for sheet_name in xl.sheet_names:
        df      = xl.parse(sheet_name)
        headers = list(df.columns)
        rows    = [[str(v) for v in row] for _, row in df.iterrows()]
        chunks.extend(_chunk_table_rows(
            headers, rows, source, "xlsx",
            extra_meta={"sheet": sheet_name},
        ))
    return chunks


def parse_pdf(data: bytes, source: str) -> List[Chunk]:
    """텍스트 + 표 분리 처리 (chromadb_chunking.PdfChunker 동일 로직)"""
    chunks = []
    with pdfplumber.open(io.BytesIO(data)) as pdf:
        for page_num, page in enumerate(pdf.pages, 1):
            # 표 처리
            tables      = page.extract_tables()
            table_bboxes = [t.bbox for t in page.find_tables()] if tables else []
            for t_idx, table in enumerate(tables):
                if not table:
                    continue
                headers = [str(h).strip() if h else f"col{i}" for i, h in enumerate(table[0])]
                rows    = [[str(c).strip() if c else "" for c in row] for row in table[1:]]
                chunks.extend(_chunk_table_rows(
                    headers, rows, source, "pdf_table",
                    rows_per_chunk=PDF_TABLE_ROWS_PER_CHUNK,
                    extra_meta={"page": page_num, "table_index": t_idx},
                ))

            # 텍스트 처리 (표 영역 제외)
            if table_bboxes:
                page_text = page.filter(
                    lambda obj: obj["object_type"] == "char"
                    and not any(
                        bbox[0] <= obj["x0"] <= bbox[2]
                        and bbox[1] <= obj["top"] <= bbox[3]
                        for bbox in table_bboxes
                    )
                ).extract_text()
            else:
                page_text = page.extract_text()

            if page_text and page_text.strip():
                chunks.extend(_chunk_text(
                    page_text.strip(), source, "pdf_text",
                    extra_meta={"page": page_num},
                ))
    return chunks


def parse_docx(data: bytes, source: str) -> List[Chunk]:
    """단락 + 표 분리 처리 (chromadb_chunking.DocxChunker 동일 로직)"""
    doc    = docx.Document(io.BytesIO(data))
    chunks = []

    # 표 처리
    for t_idx, table in enumerate(doc.tables):
        rows = [[cell.text.strip() for cell in row.cells] for row in table.rows]
        if not rows:
            continue
        headers = rows[0]
        chunks.extend(_chunk_table_rows(
            headers, rows[1:], source, "docx_table",
            extra_meta={"table_index": t_idx},
        ))

    # 단락 처리: 헤딩 기준 섹션 분리
    sections: List[Dict] = []
    cur_header, cur_lines = "No Header", []
    for para in doc.paragraphs:
        text  = para.text.strip()
        style = para.style.name
        if not text:
            continue
        if style.startswith("Heading"):
            if cur_lines:
                sections.append({"header": cur_header, "text": "\n".join(cur_lines)})
            cur_header, cur_lines = text, []
        else:
            cur_lines.append(text)
    if cur_lines:
        sections.append({"header": cur_header, "text": "\n".join(cur_lines)})

    for sec in sections:
        chunks.extend(_chunk_text(
            sec["text"], source, "docx_text",
            extra_meta={"section_header": sec["header"]},
        ))
    return chunks


def parse_pptx(data: bytes, source: str) -> List[Chunk]:
    """슬라이드별 텍스트 + 표 + 노트 처리 (chromadb_chunking.PptxChunker 동일 로직)"""
    prs    = Presentation(io.BytesIO(data))
    chunks = []

    for slide_num, slide in enumerate(prs.slides, 1):
        slide_title = (
            slide.shapes.title.text.strip()
            if slide.shapes.title and slide.shapes.title.has_text_frame
            else "No Title"
        )
        text_parts = []

        for shape in slide.shapes:
            if shape.has_table:
                table   = shape.table
                rows    = [[cell.text.strip() for cell in row.cells] for row in table.rows]
                if not rows:
                    continue
                headers = rows[0]
                chunks.extend(_chunk_table_rows(
                    headers, rows[1:], source, "pptx_table",
                    extra_meta={"slide": slide_num, "slide_title": slide_title},
                ))
            elif shape.has_text_frame:
                for para in shape.text_frame.paragraphs:
                    text = para.text.strip()
                    if text:
                        text_parts.append(text)

        if text_parts:
            chunks.extend(_chunk_text(
                "\n".join(text_parts), source, "pptx_text",
                extra_meta={"slide": slide_num, "slide_title": slide_title},
            ))

        # 슬라이드 노트
        if slide.has_notes_slide:
            notes = slide.notes_slide.notes_text_frame.text.strip()
            if notes:
                chunks.append(Chunk(
                    text=notes,
                    metadata={
                        "source":      source,
                        "doc_type":    "pptx_notes",
                        "slide":       slide_num,
                        "slide_title": slide_title,
                    },
                ))
    return chunks


# ──────────────────────────────────────────────
# 파서 라우터
# ──────────────────────────────────────────────
PARSERS = {
    "txt":  parse_txt,
    "md":   parse_md,
    "html": parse_html,
    "htm":  parse_html,   # .htm 도 동일하게 처리
    "json": parse_json,
    "csv":  parse_csv,
    "xlsx": parse_xlsx,
    "pdf":  parse_pdf,
    "docx": parse_docx,
    "pptx": parse_pptx,
}


def _get_ext(key: str) -> str:
    return key.rsplit(".", 1)[-1].lower() if "." in key else ""


def _route_and_parse(data: bytes, key: str) -> List[Chunk]:
    ext = _get_ext(key)
    if ext not in PARSERS:
        raise ValueError(f"지원하지 않는 파일 형식: .{ext} (지원: {list(PARSERS.keys())})")
    logger.info(f"[Router] 파일: {key}, 형식: {ext}")
    return PARSERS[ext](data, key)


# ──────────────────────────────────────────────
# Lambda 핸들러
# ──────────────────────────────────────────────
def lambda_handler(event: Dict, context: Any) -> Dict:
    """
    Bedrock 이 전달하는 event 구조:
    {
        "version": "1.0",
        "knowledgeBaseId": "...",
        "dataSourceId": "...",
        "ingestionJobId": "...",
        "bucketName": "...",
        "inputFiles": [
            {
                "originalFileLocation": {
                    "type": "S3",
                    "s3_location": { "uri": "s3://bucket/users/u1/raw/file.xlsx" }
                },
                "fileMetadata": {},
                "contentBatches": [{ "key": "users/u1/raw/file.xlsx" }]
            }
        ]
    }
    """
    logger.info(f"[Lambda] 이벤트 수신: {json.dumps(event)}")
    output_files = []

    for input_file in event.get("inputFiles", []):
        s3_uri   = input_file["originalFileLocation"]["s3_location"]["uri"]
        bucket, key = _parse_s3_uri(s3_uri)

        try:
            # 1. S3 raw/ 에서 파일 읽기
            data = _read_s3_bytes(bucket, key)
            logger.info(f"[Lambda] S3 읽기 완료: {s3_uri} ({len(data):,} bytes)")

            # 2. 파일 형식별 파싱 + 청킹
            chunks = _route_and_parse(data, key)
            logger.info(f"[Lambda] 청킹 완료: {len(chunks)}개 청크")

            # 3. processed/ 에 저장
            output_uri = _write_chunks_to_s3(INTERMEDIATE_BUCKET, key, chunks)

            output_files.append({
                "originalFileLocation": input_file["originalFileLocation"],
                "fileMetadata":         input_file.get("fileMetadata", {}),
                "contentBatches":       [{"key": output_uri.replace(f"s3://{INTERMEDIATE_BUCKET}/", "")}],
            })

        except ValueError as e:
            # 지원하지 않는 파일 형식 → 해당 파일만 스킵, Job 은 계속
            logger.warning(f"[Lambda] 스킵: {s3_uri} — {e}")
        except Exception as e:
            logger.error(f"[Lambda] 처리 실패: {s3_uri} — {e}", exc_info=True)
            raise  # Bedrock 이 Ingestion Job 을 FAILED 로 마킹

    logger.info(f"[Lambda] 완료: {len(output_files)}개 파일 처리")
    return {"outputFiles": output_files}